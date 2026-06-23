#!/usr/bin/env python3
"""Generate the red-hat.pptx template with 14 slide layouts.

Each layout has named placeholders with proper formatting (fonts, colors,
bullet styles, spacing). The template is designed to be loaded by
PresentationBuilder and filled with content via placeholders.

Run once to bootstrap the template, then refine in PowerPoint/LibreOffice.
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Red Hat palette
PALETTE = {
    "primary": "EE0000",
    "secondary": "151515",
    "accent": "A30000",
    "background": "FFFFFF",
    "background_alt": "F2F2F2",
    "text": "151515",
    "text_light": "FFFFFF",
    "heading": "151515",
    "subtitle": "6A6E73",
    "divider": "EE0000",
}

FONT = "Calibri"

# OpenXML namespaces
nsmap = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def pct(val: float, dimension: int) -> int:
    """Convert percentage to EMU."""
    return int(val / 100 * dimension)


def make_solid_fill(color: str) -> etree._Element:
    """Create a solidFill element."""
    fill = etree.SubElement(etree.Element("dummy"), f"{{{nsmap['a']}}}solidFill")
    etree.SubElement(fill, f"{{{nsmap['a']}}}srgbClr", val=color)
    return fill


def make_text_body_props(
    font_size: int,
    color: str,
    bold: bool = False,
    align: str = "l",
    is_bullet: bool = False,
    font_family: str = FONT,
) -> tuple[etree._Element, etree._Element]:
    """Create txBody with a default paragraph style.

    Returns (txBody element, the paragraph element inside it).
    """
    ns_a = nsmap["a"]

    txBody = etree.Element(f"{{{ns_a}}}txBody")
    etree.SubElement(txBody, f"{{{ns_a}}}bodyPr", wrap="square")
    lstStyle = etree.SubElement(txBody, f"{{{ns_a}}}lstStyle")

    defPPr = etree.SubElement(lstStyle, f"{{{ns_a}}}lvl1pPr", algn=align)

    if is_bullet:
        defPPr.set("marL", str(Inches(0.35)))
        defPPr.set("indent", str(-Inches(0.25)))
        etree.SubElement(defPPr, f"{{{ns_a}}}buFont", typeface="Arial", charset="0")
        etree.SubElement(defPPr, f"{{{ns_a}}}buChar", char="•")
    else:
        etree.SubElement(defPPr, f"{{{ns_a}}}buNone")

    defRPr = etree.SubElement(
        defPPr,
        f"{{{ns_a}}}defRPr",
        lang="en-US",
        sz=str(font_size * 100),
        b="1" if bold else "0",
    )
    fill = etree.SubElement(defRPr, f"{{{ns_a}}}solidFill")
    etree.SubElement(fill, f"{{{ns_a}}}srgbClr", val=color)
    etree.SubElement(defRPr, f"{{{ns_a}}}latin", typeface=font_family)

    p = etree.SubElement(txBody, f"{{{ns_a}}}p")
    r = etree.SubElement(p, f"{{{ns_a}}}r")
    etree.SubElement(r, f"{{{ns_a}}}rPr", lang="en-US")
    t = etree.SubElement(r, f"{{{ns_a}}}t")
    t.text = ""

    return txBody, p


def make_placeholder_sp(
    name: str,
    ph_idx: int,
    ph_type: str,
    x: int,
    y: int,
    w: int,
    h: int,
    font_size: int,
    color: str,
    bold: bool = False,
    align: str = "l",
    is_bullet: bool = False,
) -> etree._Element:
    """Create a placeholder shape element for a slide layout."""
    ns_a = nsmap["a"]
    ns_p = nsmap["p"]

    sp = etree.Element(f"{{{ns_p}}}sp")

    # Non-visual properties
    nvSpPr = etree.SubElement(sp, f"{{{ns_p}}}nvSpPr")
    etree.SubElement(nvSpPr, f"{{{ns_p}}}cNvPr", id=str(ph_idx + 2), name=name)
    cNvSpPr = etree.SubElement(nvSpPr, f"{{{ns_p}}}cNvSpPr")
    etree.SubElement(cNvSpPr, f"{{{ns_a}}}spLocks", noGrp="1")
    nvPr = etree.SubElement(nvSpPr, f"{{{ns_p}}}nvPr")
    etree.SubElement(nvPr, f"{{{ns_p}}}ph", type=ph_type, idx=str(ph_idx))

    # Shape properties (position and size)
    spPr = etree.SubElement(sp, f"{{{ns_p}}}spPr")
    xfrm = etree.SubElement(spPr, f"{{{ns_a}}}xfrm")
    etree.SubElement(xfrm, f"{{{ns_a}}}off", x=str(x), y=str(y))
    etree.SubElement(xfrm, f"{{{ns_a}}}ext", cx=str(w), cy=str(h))

    # Text body with default formatting
    txBody, _ = make_text_body_props(font_size, color, bold, align, is_bullet)
    sp.append(txBody)

    return sp


def make_bg_fill(color: str) -> etree._Element:
    """Create a background element with solid fill."""
    ns_p = nsmap["p"]
    ns_a = nsmap["a"]

    bg = etree.Element(f"{{{ns_p}}}bg")
    bgPr = etree.SubElement(bg, f"{{{ns_p}}}bgPr")
    fill = etree.SubElement(bgPr, f"{{{ns_a}}}solidFill")
    etree.SubElement(fill, f"{{{ns_a}}}srgbClr", val=color)
    etree.SubElement(bgPr, f"{{{ns_a}}}effectLst")

    return bg


def make_accent_bar(y_pct: float, w_pct: float, color: str = PALETTE["divider"]) -> etree._Element:
    """Create a horizontal accent bar shape."""
    ns_a = nsmap["a"]
    ns_p = nsmap["p"]

    sp = etree.Element(f"{{{ns_p}}}sp")
    nvSpPr = etree.SubElement(sp, f"{{{ns_p}}}nvSpPr")
    etree.SubElement(nvSpPr, f"{{{ns_p}}}cNvPr", id="100", name="Accent Bar")
    etree.SubElement(nvSpPr, f"{{{ns_p}}}cNvSpPr")
    etree.SubElement(nvSpPr, f"{{{ns_p}}}nvPr")

    spPr = etree.SubElement(sp, f"{{{ns_p}}}spPr")
    xfrm = etree.SubElement(spPr, f"{{{ns_a}}}xfrm")
    etree.SubElement(xfrm, f"{{{ns_a}}}off", x="0", y=str(pct(y_pct, SLIDE_H)))
    etree.SubElement(xfrm, f"{{{ns_a}}}ext", cx=str(pct(w_pct, SLIDE_W)), cy=str(Pt(4)))
    etree.SubElement(spPr, f"{{{ns_a}}}prstGeom", prst="rect")
    fill = etree.SubElement(spPr, f"{{{ns_a}}}solidFill")
    etree.SubElement(fill, f"{{{ns_a}}}srgbClr", val=color)
    ln = etree.SubElement(spPr, f"{{{ns_a}}}ln")
    etree.SubElement(ln, f"{{{ns_a}}}noFill")

    return sp


def make_slide_number_sp() -> etree._Element:
    """Create a slide number placeholder in the bottom-right corner."""
    ns_a = nsmap["a"]
    ns_p = nsmap["p"]

    sp = etree.Element(f"{{{ns_p}}}sp")
    nvSpPr = etree.SubElement(sp, f"{{{ns_p}}}nvSpPr")
    etree.SubElement(nvSpPr, f"{{{ns_p}}}cNvPr", id="101", name="Slide Number")
    cNvSpPr = etree.SubElement(nvSpPr, f"{{{ns_p}}}cNvSpPr")
    etree.SubElement(cNvSpPr, f"{{{ns_a}}}spLocks", noGrp="1")
    nvPr = etree.SubElement(nvSpPr, f"{{{ns_p}}}nvPr")
    etree.SubElement(nvPr, f"{{{ns_p}}}ph", type="sldNum", sz="quarter", idx="20")

    spPr = etree.SubElement(sp, f"{{{ns_p}}}spPr")
    xfrm = etree.SubElement(spPr, f"{{{ns_a}}}xfrm")
    etree.SubElement(xfrm, f"{{{ns_a}}}off", x=str(pct(92.5, SLIDE_W)), y=str(pct(93.3, SLIDE_H)))
    etree.SubElement(xfrm, f"{{{ns_a}}}ext", cx=str(Inches(0.75)), cy=str(Inches(0.35)))

    txBody = etree.SubElement(sp, f"{{{ns_p}}}txBody")
    etree.SubElement(txBody, f"{{{ns_a}}}bodyPr")
    etree.SubElement(txBody, f"{{{ns_a}}}lstStyle")
    p = etree.SubElement(txBody, f"{{{ns_a}}}p")
    etree.SubElement(p, f"{{{ns_a}}}pPr", algn="r")
    fld = etree.SubElement(p, f"{{{ns_a}}}fld", type="slidenum")
    fld.set("id", "{B6F15528-F159-4107-8E1A-4C7D40CC3D3E}")
    rPr = etree.SubElement(fld, f"{{{ns_a}}}rPr", lang="en-US", sz="1400")
    fill = etree.SubElement(rPr, f"{{{ns_a}}}solidFill")
    etree.SubElement(fill, f"{{{ns_a}}}srgbClr", val=PALETTE["subtitle"])
    t = etree.SubElement(fld, f"{{{ns_a}}}t")
    t.text = "‹#›"

    return sp


# Layout definitions: name -> list of placeholders
# Each placeholder: (name, idx, type, x%, y%, w%, h%, font_size, color, bold, align, is_bullet)
LAYOUTS = {
    "title": {
        "bg": PALETTE["background"],
        "accent_bar": (48.0, 22.5),
        "slide_number": False,
        "placeholders": [
            (
                "Title",
                0,
                "ctrTitle",
                7.5,
                26.7,
                85.0,
                20.0,
                28,
                PALETTE["heading"],
                True,
                "l",
                False,
            ),
            (
                "Subtitle",
                1,
                "subTitle",
                7.5,
                53.3,
                85.0,
                13.3,
                24,
                PALETTE["subtitle"],
                False,
                "l",
                False,
            ),
        ],
    },
    "title-dark": {
        "bg": PALETTE["background_alt"],
        "accent_bar": (48.0, 22.5),
        "slide_number": False,
        "placeholders": [
            (
                "Title",
                0,
                "ctrTitle",
                7.5,
                26.7,
                85.0,
                20.0,
                28,
                PALETTE["heading"],
                True,
                "l",
                False,
            ),
            (
                "Subtitle",
                1,
                "subTitle",
                7.5,
                53.3,
                85.0,
                13.3,
                24,
                PALETTE["subtitle"],
                False,
                "l",
                False,
            ),
        ],
    },
    "section": {
        "bg": PALETTE["background_alt"],
        "accent_bar": (54.7, 22.5),
        "slide_number": True,
        "placeholders": [
            ("Title", 0, "title", 11.3, 33.3, 77.5, 20.0, 28, PALETTE["heading"], True, "l", False),
            (
                "Subtitle",
                1,
                "body",
                11.3,
                60.0,
                77.5,
                10.7,
                24,
                PALETTE["subtitle"],
                False,
                "l",
                False,
            ),
        ],
    },
    "content": {
        "bg": PALETTE["background"],
        "accent_bar": (16.7, 15.0),
        "slide_number": True,
        "placeholders": [
            ("Title", 0, "title", 5.6, 5.3, 88.8, 10.7, 28, PALETTE["heading"], True, "l", False),
            ("Body", 1, "body", 5.6, 21.3, 88.8, 68.7, 20, PALETTE["text"], False, "l", True),
        ],
    },
    "content-with-icon": {
        "bg": PALETTE["background"],
        "accent_bar": (16.7, 15.0),
        "slide_number": True,
        "placeholders": [
            ("Title", 0, "title", 5.6, 5.3, 88.8, 10.7, 28, PALETTE["heading"], True, "l", False),
            ("Body", 1, "body", 5.6, 21.3, 58.0, 68.7, 20, PALETTE["text"], False, "l", True),
            ("Icon", 10, "pic", 68.0, 25.0, 25.0, 33.3, 18, PALETTE["text"], False, "l", False),
        ],
    },
    "content-with-graphic": {
        "bg": PALETTE["background"],
        "accent_bar": (16.7, 15.0),
        "slide_number": True,
        "placeholders": [
            ("Title", 0, "title", 5.6, 5.3, 88.8, 10.7, 28, PALETTE["heading"], True, "l", False),
            ("Body", 1, "body", 5.6, 21.3, 48.0, 68.7, 20, PALETTE["text"], False, "l", True),
            ("Graphic", 10, "pic", 56.0, 21.3, 38.0, 68.7, 18, PALETTE["text"], False, "l", False),
        ],
    },
    "two-column": {
        "bg": PALETTE["background"],
        "accent_bar": (16.7, 15.0),
        "slide_number": True,
        "placeholders": [
            ("Title", 0, "title", 5.6, 5.3, 88.8, 10.7, 28, PALETTE["heading"], True, "l", False),
            (
                "Left Heading",
                10,
                "body",
                5.6,
                21.3,
                42.0,
                8.0,
                20,
                PALETTE["primary"],
                True,
                "l",
                False,
            ),
            ("Left Body", 11, "body", 5.6, 30.7, 42.0, 59.3, 20, PALETTE["text"], False, "l", True),
            (
                "Right Heading",
                12,
                "body",
                52.4,
                21.3,
                42.0,
                8.0,
                20,
                PALETTE["primary"],
                True,
                "l",
                False,
            ),
            (
                "Right Body",
                13,
                "body",
                52.4,
                30.7,
                42.0,
                59.3,
                20,
                PALETTE["text"],
                False,
                "l",
                True,
            ),
        ],
    },
    "two-column-uneven": {
        "bg": PALETTE["background"],
        "accent_bar": (16.7, 15.0),
        "slide_number": True,
        "placeholders": [
            ("Title", 0, "title", 5.6, 5.3, 88.8, 10.7, 28, PALETTE["heading"], True, "l", False),
            (
                "Left Heading",
                10,
                "body",
                5.6,
                21.3,
                56.0,
                8.0,
                20,
                PALETTE["primary"],
                True,
                "l",
                False,
            ),
            ("Left Body", 11, "body", 5.6, 30.7, 56.0, 59.3, 20, PALETTE["text"], False, "l", True),
            (
                "Right Heading",
                12,
                "body",
                65.0,
                21.3,
                29.4,
                8.0,
                20,
                PALETTE["primary"],
                True,
                "l",
                False,
            ),
            (
                "Right Body",
                13,
                "body",
                65.0,
                30.7,
                29.4,
                59.3,
                20,
                PALETTE["text"],
                False,
                "l",
                True,
            ),
        ],
    },
    "image": {
        "bg": PALETTE["background"],
        "accent_bar": None,
        "slide_number": True,
        "placeholders": [
            ("Title", 0, "title", 5.6, 5.3, 88.8, 10.7, 28, PALETTE["heading"], True, "l", False),
            ("Image", 10, "pic", 7.5, 20.0, 85.0, 64.0, 18, PALETTE["text"], False, "l", False),
            (
                "Caption",
                11,
                "body",
                7.5,
                89.3,
                85.0,
                6.7,
                16,
                PALETTE["subtitle"],
                False,
                "ctr",
                False,
            ),
        ],
    },
    "image-with-text": {
        "bg": PALETTE["background"],
        "accent_bar": None,
        "slide_number": True,
        "placeholders": [
            ("Title", 0, "title", 5.6, 5.3, 88.8, 10.7, 28, PALETTE["heading"], True, "l", False),
            ("Image", 10, "pic", 5.6, 21.3, 45.0, 68.7, 18, PALETTE["text"], False, "l", False),
            ("Body", 1, "body", 54.0, 21.3, 40.4, 68.7, 20, PALETTE["text"], False, "l", True),
        ],
    },
    "comparison": {
        "bg": PALETTE["background"],
        "accent_bar": (16.7, 15.0),
        "slide_number": True,
        "placeholders": [
            ("Title", 0, "title", 5.6, 5.3, 88.8, 10.7, 28, PALETTE["heading"], True, "l", False),
            (
                "Left Heading",
                10,
                "body",
                5.6,
                21.3,
                42.0,
                8.0,
                20,
                PALETTE["accent"],
                True,
                "l",
                False,
            ),
            ("Left Body", 11, "body", 5.6, 30.7, 42.0, 59.3, 20, PALETTE["text"], False, "l", True),
            (
                "Right Heading",
                12,
                "body",
                52.4,
                21.3,
                42.0,
                8.0,
                20,
                PALETTE["primary"],
                True,
                "l",
                False,
            ),
            (
                "Right Body",
                13,
                "body",
                52.4,
                30.7,
                42.0,
                59.3,
                20,
                PALETTE["text"],
                False,
                "l",
                True,
            ),
        ],
    },
    "data": {
        "bg": PALETTE["background"],
        "accent_bar": None,
        "slide_number": True,
        "placeholders": [
            ("Title", 0, "title", 5.6, 5.3, 88.8, 10.7, 28, PALETTE["heading"], True, "l", False),
            (
                "Metric",
                10,
                "body",
                10.0,
                26.7,
                80.0,
                30.0,
                36,
                PALETTE["primary"],
                True,
                "ctr",
                False,
            ),
            ("Body", 1, "body", 15.0, 60.0, 70.0, 26.7, 20, PALETTE["text"], False, "ctr", True),
        ],
    },
    "quote": {
        "bg": PALETTE["background_alt"],
        "accent_bar": None,
        "slide_number": True,
        "placeholders": [
            (
                "Quote",
                0,
                "title",
                11.3,
                20.0,
                77.5,
                46.7,
                28,
                PALETTE["heading"],
                False,
                "ctr",
                False,
            ),
            (
                "Attribution",
                1,
                "body",
                11.3,
                70.0,
                77.5,
                10.0,
                16,
                PALETTE["subtitle"],
                False,
                "ctr",
                False,
            ),
        ],
    },
    "closing": {
        "bg": PALETTE["background_alt"],
        "accent_bar": None,
        "slide_number": False,
        "placeholders": [
            (
                "Title",
                0,
                "ctrTitle",
                11.3,
                26.7,
                77.5,
                20.0,
                36,
                PALETTE["heading"],
                True,
                "ctr",
                False,
            ),
            (
                "Subtitle",
                1,
                "subTitle",
                11.3,
                46.7,
                77.5,
                10.7,
                24,
                PALETTE["subtitle"],
                False,
                "ctr",
                False,
            ),
            (
                "Contact",
                10,
                "body",
                11.3,
                60.0,
                77.5,
                6.7,
                20,
                PALETTE["primary"],
                False,
                "ctr",
                False,
            ),
        ],
    },
}


def generate_template(output_path: str) -> None:
    """Generate the template .pptx with all 14 layouts."""
    # Start from default template
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Get the slide master
    master = prs.slide_masters[0]

    # Remove all existing slide layouts except the first (we'll replace it)
    # We need to work at the XML level
    master_elem = master.element

    # Remove existing layout references
    ns_p = nsmap["p"]
    ns_r = nsmap["r"]

    # Get existing layout IDs to remove
    sldLayoutIdLst = master_elem.find(f"{{{ns_p}}}sldLayoutIdLst")

    # Track which layouts to remove
    existing_layout_rIds = []
    if sldLayoutIdLst is not None:
        for sldLayoutId in sldLayoutIdLst.findall(f"{{{ns_p}}}sldLayoutId"):
            existing_layout_rIds.append(sldLayoutId.get(f"{{{ns_r}}}id"))

    # Instead of complex XML surgery, let's use python-pptx's approach:
    # Create the presentation, then for each layout, add a slide to verify it works
    # We'll build on the default template's first layout

    # Actually, the simplest approach: use the default presentation and its existing
    # layouts as a base. We'll modify layouts by rebuilding from the default.

    import io
    import zipfile

    # Save the base presentation first
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    # Read as zip, modify layout files, write new zip
    with zipfile.ZipFile(buf, "r") as zin:
        out_buf = io.BytesIO()
        with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:
            # Copy most files as-is
            skip_files = set()

            # Find existing layout files
            existing_layouts = [n for n in zin.namelist() if n.startswith("ppt/slideLayouts/")]
            skip_files.update(existing_layouts)

            # Also need to update content types and relationships
            skip_files.add("[Content_Types].xml")
            skip_files.add("ppt/slideMasters/_rels/slideMaster1.xml.rels")
            skip_files.add("ppt/slideMasters/slideMaster1.xml")

            for item in zin.namelist():
                if item not in skip_files:
                    zout.writestr(item, zin.read(item))

            # Get existing master XML
            master_xml_str = zin.read("ppt/slideMasters/slideMaster1.xml")
            master_xml = etree.fromstring(master_xml_str)

            # Get existing master rels
            master_rels_str = zin.read("ppt/slideMasters/_rels/slideMaster1.xml.rels")
            master_rels = etree.fromstring(master_rels_str)

            # Get content types
            ct_str = zin.read("[Content_Types].xml")
            ct_xml = etree.fromstring(ct_str)

            # Remove existing layout relationships from master
            ns_rels = "http://schemas.openxmlformats.org/package/2006/relationships"
            layout_rel_type = (
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
            )
            for rel in list(master_rels):
                if rel.get("Type") == layout_rel_type:
                    master_rels.remove(rel)

            # Remove existing layout ID list from master
            sldLayoutIdLst = master_xml.find(f"{{{ns_p}}}sldLayoutIdLst")
            if sldLayoutIdLst is not None:
                master_xml.remove(sldLayoutIdLst)
            sldLayoutIdLst = etree.SubElement(master_xml, f"{{{ns_p}}}sldLayoutIdLst")

            # Remove existing layout content type overrides
            for override in list(ct_xml):
                pn = override.get("PartName", "")
                if "/slideLayouts/" in pn:
                    ct_xml.remove(override)

            # Base layout XML to clone
            base_layout_xml = zin.read(existing_layouts[0])

            # Write new layouts
            for layout_idx, (layout_name, layout_def) in enumerate(LAYOUTS.items()):
                layout_num = layout_idx + 1
                layout_filename = f"ppt/slideLayouts/slideLayout{layout_num}.xml"

                # Clone base layout
                layout_xml = etree.fromstring(base_layout_xml)

                # Set name
                cSld = layout_xml.find(f"{{{ns_p}}}cSld")
                if cSld is not None:
                    cSld.set("name", layout_name)

                # Add background
                bg_color = layout_def["bg"]
                existing_bg = cSld.find(f"{{{ns_p}}}bg")
                if existing_bg is not None:
                    cSld.remove(existing_bg)
                bg_elem = make_bg_fill(bg_color)
                cSld.insert(0, bg_elem)

                # Clear existing shapes
                spTree = cSld.find(f"{{{ns_p}}}spTree")
                if spTree is not None:
                    for child in list(spTree):
                        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                        if tag == "sp":
                            spTree.remove(child)

                    # Add new placeholders
                    for ph_def in layout_def["placeholders"]:
                        (
                            name,
                            idx,
                            ph_type,
                            x_pct,
                            y_pct,
                            w_pct,
                            h_pct,
                            font_size,
                            color,
                            bold,
                            align,
                            is_bullet,
                        ) = ph_def
                        sp = make_placeholder_sp(
                            name=name,
                            ph_idx=idx,
                            ph_type=ph_type,
                            x=pct(x_pct, SLIDE_W),
                            y=pct(y_pct, SLIDE_H),
                            w=pct(w_pct, SLIDE_W),
                            h=pct(h_pct, SLIDE_H),
                            font_size=font_size,
                            color=color,
                            bold=bold,
                            align=align,
                            is_bullet=is_bullet,
                        )
                        spTree.append(sp)

                    # Add accent bar if defined
                    bar = layout_def.get("accent_bar")
                    if bar is not None:
                        bar_sp = make_accent_bar(bar[0], bar[1])
                        spTree.append(bar_sp)

                    # Add slide number placeholder if enabled
                    if layout_def.get("slide_number", False):
                        sn_sp = make_slide_number_sp()
                        spTree.append(sn_sp)

                # Write layout file
                layout_bytes = etree.tostring(
                    layout_xml, xml_declaration=True, encoding="UTF-8", standalone=True
                )
                zout.writestr(layout_filename, layout_bytes)

                # Write layout rels (relationship to master)
                rels_xml = etree.Element(f"{{{ns_rels}}}Relationships")
                etree.SubElement(
                    rels_xml,
                    f"{{{ns_rels}}}Relationship",
                    Id="rId1",
                    Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster",
                    Target="../slideMasters/slideMaster1.xml",
                )
                rels_filename = f"ppt/slideLayouts/_rels/slideLayout{layout_num}.xml.rels"
                rels_bytes = etree.tostring(
                    rels_xml,
                    xml_declaration=True,
                    encoding="UTF-8",
                    standalone=True,
                )
                zout.writestr(rels_filename, rels_bytes)

                # Add relationship in master
                rId = f"rId{layout_idx + 100}"  # Use high IDs to avoid conflicts
                etree.SubElement(
                    master_rels,
                    f"{{{ns_rels}}}Relationship",
                    Id=rId,
                    Type=layout_rel_type,
                    Target=f"../slideLayouts/slideLayout{layout_num}.xml",
                )

                # Add layout ID in master
                etree.SubElement(
                    sldLayoutIdLst,
                    f"{{{ns_p}}}sldLayoutId",
                    id=str(2147483649 + layout_idx),
                    attrib={f"{{{ns_r}}}id": rId},
                )

                # Add content type override
                etree.SubElement(
                    ct_xml,
                    "Override",
                    PartName=f"/{layout_filename}",
                    ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml",
                )

            # Write modified master
            master_bytes = etree.tostring(
                master_xml,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )
            zout.writestr("ppt/slideMasters/slideMaster1.xml", master_bytes)

            # Write modified master rels
            master_rels_bytes = etree.tostring(
                master_rels,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )
            zout.writestr(
                "ppt/slideMasters/_rels/slideMaster1.xml.rels",
                master_rels_bytes,
            )

            # Write modified content types
            ct_bytes = etree.tostring(
                ct_xml, xml_declaration=True, encoding="UTF-8", standalone=True
            )
            zout.writestr("[Content_Types].xml", ct_bytes)

    # Save output
    out_path = Path(output_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_bytes(out_buf.getvalue())
    print(f"✓ Generated template: {out_path}")
    print(f"  Layouts: {len(LAYOUTS)}")
    for i, name in enumerate(LAYOUTS):
        print(f"    [{i}] {name}")


if __name__ == "__main__":
    output = Path(__file__).parent / "default.pptx"
    generate_template(str(output))
