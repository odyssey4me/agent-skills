#!/usr/bin/env python3
"""Google Slides integration skill for AI agents.

Markdown-driven presentation builder with Google Slides upload support.
Write presentations in Markdown, build .pptx locally, review in LibreOffice,
then upload to Google Slides.

Usage:
    python google-slides.py check
    python google-slides.py auth setup --client-id ID --client-secret SECRET
    python google-slides.py create --file deck.md --output deck.pptx
    python google-slides.py create --title "My Deck" --file deck.md
    python google-slides.py get PRESENTATION_ID
    python google-slides.py get PRESENTATION_ID -o deck.md
    python google-slides.py update PRESENTATION_ID --file deck.pptx
    python google-slides.py palettes

Requirements:
    pip install --user google-auth google-auth-oauthlib google-api-python-client keyring pyyaml
    pip install --user python-pptx cairosvg  (for local presentation building)
"""

from __future__ import annotations

# Standard library imports
import argparse
import contextlib
import json
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any

# ============================================================================
# DEPENDENCY CHECKS
# ============================================================================

try:
    from google.auth.transport.requests import Request
    from google.oauth2.credentials import Credentials
    from google_auth_oauthlib.flow import InstalledAppFlow

    GOOGLE_AUTH_AVAILABLE = True
except ImportError:
    GOOGLE_AUTH_AVAILABLE = False

try:
    from googleapiclient.discovery import build
    from googleapiclient.errors import HttpError

    GOOGLE_API_CLIENT_AVAILABLE = True
except ImportError:
    GOOGLE_API_CLIENT_AVAILABLE = False

try:
    import keyring

    KEYRING_AVAILABLE = True
except ImportError:
    KEYRING_AVAILABLE = False

try:
    import yaml

    YAML_AVAILABLE = True
except ImportError:
    YAML_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import cairosvg  # noqa: F401

    CAIROSVG_AVAILABLE = True
except ImportError:
    CAIROSVG_AVAILABLE = False


# ============================================================================
# CONSTANTS
# ============================================================================

SERVICE_NAME = "agent-skills"
CONFIG_DIR = Path.home() / ".config" / "agent-skills"
ICON_CACHE_DIR = Path.home() / ".cache" / "agent-skills" / "icons"

# Google Slides API scopes
SLIDES_SCOPES_READONLY = ["https://www.googleapis.com/auth/presentations.readonly"]
SLIDES_SCOPES = ["https://www.googleapis.com/auth/presentations"]
SLIDES_SCOPES_DEFAULT = SLIDES_SCOPES_READONLY

# Drive API scopes
DRIVE_SCOPES_READONLY = ["https://www.googleapis.com/auth/drive.readonly"]
DRIVE_SCOPES_WRITE = ["https://www.googleapis.com/auth/drive.file"]

# EMU (English Metric Units) conversion
EMU_PER_INCH = 914400
EMU_PER_PT = 12700

# Slide dimensions (in EMU via Inches helper)
SLIDE_WIDTH_WIDESCREEN = Inches(13.333) if PPTX_AVAILABLE else 12192000
SLIDE_HEIGHT_WIDESCREEN = Inches(7.5) if PPTX_AVAILABLE else 6858000
SLIDE_WIDTH_STANDARD = Inches(10) if PPTX_AVAILABLE else 9144000
SLIDE_HEIGHT_STANDARD = Inches(7.5) if PPTX_AVAILABLE else 6858000


# ============================================================================
# COLOR PALETTES
# ============================================================================

COLOR_PALETTES: dict[str, dict[str, str]] = {
    "red-hat": {
        "primary": "#EE0000",
        "secondary": "#151515",
        "accent": "#A30000",
        "background": "#FFFFFF",
        "background_alt": "#F2F2F2",
        "text": "#151515",
        "text_light": "#FFFFFF",
        "heading": "#151515",
        "subtitle": "#6A6E73",
        "divider": "#EE0000",
    },
    "red-hat-dark": {
        "primary": "#EE0000",
        "secondary": "#FFFFFF",
        "accent": "#F56E6E",
        "background": "#151515",
        "background_alt": "#1F1F1F",
        "text": "#E0E0E0",
        "text_light": "#FFFFFF",
        "heading": "#FFFFFF",
        "subtitle": "#A3A3A3",
        "divider": "#EE0000",
    },
    "red-hat-teal": {
        "primary": "#37A3A3",
        "secondary": "#151515",
        "accent": "#EE0000",
        "background": "#FFFFFF",
        "background_alt": "#DAF2F2",
        "text": "#151515",
        "text_light": "#FFFFFF",
        "heading": "#004D4D",
        "subtitle": "#6A6E73",
        "divider": "#37A3A3",
    },
    "red-hat-purple": {
        "primary": "#5E40BE",
        "secondary": "#151515",
        "accent": "#EE0000",
        "background": "#FFFFFF",
        "background_alt": "#ECE6FF",
        "text": "#151515",
        "text_light": "#FFFFFF",
        "heading": "#21134D",
        "subtitle": "#6A6E73",
        "divider": "#5E40BE",
    },
    "corporate-blue": {
        "primary": "#0066CC",
        "secondary": "#003366",
        "accent": "#0099FF",
        "background": "#FFFFFF",
        "background_alt": "#F0F4F8",
        "text": "#1A1A1A",
        "text_light": "#FFFFFF",
        "heading": "#003366",
        "subtitle": "#666666",
        "divider": "#0066CC",
    },
    "dark-mode": {
        "primary": "#BB86FC",
        "secondary": "#03DAC6",
        "accent": "#CF6679",
        "background": "#121212",
        "background_alt": "#1E1E1E",
        "text": "#E0E0E0",
        "text_light": "#FFFFFF",
        "heading": "#FFFFFF",
        "subtitle": "#B0B0B0",
        "divider": "#BB86FC",
    },
    "minimal-light": {
        "primary": "#333333",
        "secondary": "#666666",
        "accent": "#0066CC",
        "background": "#FFFFFF",
        "background_alt": "#FAFAFA",
        "text": "#333333",
        "text_light": "#FFFFFF",
        "heading": "#1A1A1A",
        "subtitle": "#888888",
        "divider": "#DDDDDD",
    },
}

FONT_CONFIG = {
    "title_size": 36,
    "subtitle_size": 24,
    "heading_size": 28,
    "body_size": 20,
    "caption_size": 16,
    "font_family": "Calibri",
}

# ============================================================================
# SLIDE LAYOUT TEMPLATES
# ============================================================================
# Positions as percentages of slide dimensions. Converted to EMU at build time.
# Derived from existing hard-coded Inches() values on 13.333x7.5" widescreen.
#
# Placeholder roles:
#   "text"   — rendered via _add_text_box()
#   "bullets" — rendered via _add_bullet_list()
#   "image"  — rendered via slide.shapes.add_picture()

SLIDE_LAYOUTS: dict[str, dict[str, Any]] = {
    "title": {
        "background": "background",
        "accent_bar": {"y": 48.0, "w": 22.5},
        "slide_number": False,
        "placeholders": {
            "title": {
                "x": 7.5,
                "y": 26.7,
                "w": 85.0,
                "h": 20.0,
                "role": "text",
                "font": "title_size",
                "color": "heading",
                "bold": True,
                "align": "left",
            },
            "subtitle": {
                "x": 7.5,
                "y": 53.3,
                "w": 85.0,
                "h": 13.3,
                "role": "text",
                "font": "subtitle_size",
                "color": "subtitle",
            },
        },
    },
    "title-dark": {
        "background": "background_alt",
        "accent_bar": {"y": 48.0, "w": 22.5},
        "slide_number": False,
        "placeholders": {
            "title": {
                "x": 7.5,
                "y": 26.7,
                "w": 85.0,
                "h": 20.0,
                "role": "text",
                "font": "title_size",
                "color": "heading",
                "bold": True,
                "align": "left",
            },
            "subtitle": {
                "x": 7.5,
                "y": 53.3,
                "w": 85.0,
                "h": 8.0,
                "role": "text",
                "font": "subtitle_size",
                "color": "subtitle",
            },
        },
    },
    "section": {
        "background": "background_alt",
        "accent_bar": {"y": 54.7, "w": 22.5},
        "slide_number": True,
        "placeholders": {
            "title": {
                "x": 11.3,
                "y": 33.3,
                "w": 77.5,
                "h": 20.0,
                "role": "text",
                "font": "title_size",
                "color": "heading",
                "bold": True,
                "align": "left",
            },
            "subtitle": {
                "x": 11.3,
                "y": 60.0,
                "w": 77.5,
                "h": 10.7,
                "role": "text",
                "font": "subtitle_size",
                "color": "subtitle",
            },
        },
    },
    "content": {
        "background": "background",
        "accent_bar": {"y": 16.7, "w": 15.0},
        "slide_number": True,
        "placeholders": {
            "title": {
                "x": 5.6,
                "y": 5.3,
                "w": 88.8,
                "h": 10.7,
                "role": "text",
                "font": "heading_size",
                "color": "heading",
                "bold": True,
            },
            "body": {
                "x": 5.6,
                "y": 21.3,
                "w": 88.8,
                "h": 68.7,
                "role": "bullets",
                "font": "body_size",
                "color": "text",
            },
        },
    },
    "content-with-icon": {
        "background": "background",
        "accent_bar": {"y": 16.7, "w": 15.0},
        "slide_number": True,
        "placeholders": {
            "title": {
                "x": 5.6,
                "y": 5.3,
                "w": 88.8,
                "h": 10.7,
                "role": "text",
                "font": "heading_size",
                "color": "heading",
                "bold": True,
            },
            "body": {
                "x": 5.6,
                "y": 21.3,
                "w": 58.0,
                "h": 68.7,
                "role": "bullets",
                "font": "body_size",
                "color": "text",
            },
            "icon": {
                "x": 68.0,
                "y": 25.0,
                "w": 25.0,
                "h": 33.3,
                "role": "image",
            },
        },
    },
    "content-with-graphic": {
        "background": "background",
        "accent_bar": {"y": 16.7, "w": 15.0},
        "slide_number": True,
        "placeholders": {
            "title": {
                "x": 5.6,
                "y": 5.3,
                "w": 88.8,
                "h": 10.7,
                "role": "text",
                "font": "heading_size",
                "color": "heading",
                "bold": True,
            },
            "body": {
                "x": 5.6,
                "y": 21.3,
                "w": 48.0,
                "h": 68.7,
                "role": "bullets",
                "font": "body_size",
                "color": "text",
            },
            "graphic": {
                "x": 56.0,
                "y": 21.3,
                "w": 38.0,
                "h": 68.7,
                "role": "image",
            },
        },
    },
    "two-column": {
        "background": "background",
        "accent_bar": {"y": 16.7, "w": 15.0},
        "slide_number": True,
        "placeholders": {
            "title": {
                "x": 5.6,
                "y": 5.3,
                "w": 88.8,
                "h": 10.7,
                "role": "text",
                "font": "heading_size",
                "color": "heading",
                "bold": True,
            },
            "left_heading": {
                "x": 5.6,
                "y": 21.3,
                "w": 42.0,
                "h": 8.0,
                "role": "text",
                "font": "subtitle_size",
                "color": "primary",
                "bold": True,
            },
            "left_body": {
                "x": 5.6,
                "y": 30.7,
                "w": 42.0,
                "h": 59.3,
                "role": "bullets",
                "font": "body_size",
                "color": "text",
            },
            "right_heading": {
                "x": 52.4,
                "y": 21.3,
                "w": 42.0,
                "h": 8.0,
                "role": "text",
                "font": "subtitle_size",
                "color": "primary",
                "bold": True,
            },
            "right_body": {
                "x": 52.4,
                "y": 30.7,
                "w": 42.0,
                "h": 59.3,
                "role": "bullets",
                "font": "body_size",
                "color": "text",
            },
        },
    },
    "two-column-uneven": {
        "background": "background",
        "accent_bar": {"y": 16.7, "w": 15.0},
        "slide_number": True,
        "placeholders": {
            "title": {
                "x": 5.6,
                "y": 5.3,
                "w": 88.8,
                "h": 10.7,
                "role": "text",
                "font": "heading_size",
                "color": "heading",
                "bold": True,
            },
            "left_heading": {
                "x": 5.6,
                "y": 21.3,
                "w": 56.0,
                "h": 8.0,
                "role": "text",
                "font": "subtitle_size",
                "color": "primary",
                "bold": True,
            },
            "left_body": {
                "x": 5.6,
                "y": 30.7,
                "w": 56.0,
                "h": 59.3,
                "role": "bullets",
                "font": "body_size",
                "color": "text",
            },
            "right_heading": {
                "x": 65.0,
                "y": 21.3,
                "w": 29.4,
                "h": 8.0,
                "role": "text",
                "font": "subtitle_size",
                "color": "primary",
                "bold": True,
            },
            "right_body": {
                "x": 65.0,
                "y": 30.7,
                "w": 29.4,
                "h": 59.3,
                "role": "bullets",
                "font": "body_size",
                "color": "text",
            },
        },
    },
    "image": {
        "background": "background",
        "accent_bar": None,
        "slide_number": True,
        "placeholders": {
            "title": {
                "x": 5.6,
                "y": 5.3,
                "w": 88.8,
                "h": 10.7,
                "role": "text",
                "font": "heading_size",
                "color": "heading",
                "bold": True,
            },
            "image": {
                "x": 7.5,
                "y": 20.0,
                "w": 85.0,
                "h": 64.0,
                "role": "image",
            },
            "caption": {
                "x": 7.5,
                "y": 89.3,
                "w": 85.0,
                "h": 6.7,
                "role": "text",
                "font": "caption_size",
                "color": "subtitle",
                "align": "center",
            },
        },
    },
    "image-with-text": {
        "background": "background",
        "accent_bar": None,
        "slide_number": True,
        "placeholders": {
            "title": {
                "x": 5.6,
                "y": 5.3,
                "w": 88.8,
                "h": 10.7,
                "role": "text",
                "font": "heading_size",
                "color": "heading",
                "bold": True,
            },
            "image": {
                "x": 5.6,
                "y": 21.3,
                "w": 45.0,
                "h": 68.7,
                "role": "image",
            },
            "body": {
                "x": 54.0,
                "y": 21.3,
                "w": 40.4,
                "h": 68.7,
                "role": "bullets",
                "font": "body_size",
                "color": "text",
            },
        },
    },
    "comparison": {
        "background": "background",
        "accent_bar": {"y": 16.7, "w": 15.0},
        "slide_number": True,
        "placeholders": {
            "title": {
                "x": 5.6,
                "y": 5.3,
                "w": 88.8,
                "h": 10.7,
                "role": "text",
                "font": "heading_size",
                "color": "heading",
                "bold": True,
            },
            "left_heading": {
                "x": 5.6,
                "y": 21.3,
                "w": 42.0,
                "h": 8.0,
                "role": "text",
                "font": "subtitle_size",
                "color": "accent",
                "bold": True,
            },
            "left_body": {
                "x": 5.6,
                "y": 30.7,
                "w": 42.0,
                "h": 59.3,
                "role": "bullets",
                "font": "body_size",
                "color": "text",
            },
            "right_heading": {
                "x": 52.4,
                "y": 21.3,
                "w": 42.0,
                "h": 8.0,
                "role": "text",
                "font": "subtitle_size",
                "color": "primary",
                "bold": True,
            },
            "right_body": {
                "x": 52.4,
                "y": 30.7,
                "w": 42.0,
                "h": 59.3,
                "role": "bullets",
                "font": "body_size",
                "color": "text",
            },
        },
    },
    "data": {
        "background": "background",
        "accent_bar": None,
        "slide_number": True,
        "placeholders": {
            "title": {
                "x": 5.6,
                "y": 5.3,
                "w": 88.8,
                "h": 10.7,
                "role": "text",
                "font": "heading_size",
                "color": "heading",
                "bold": True,
            },
            "metric": {
                "x": 10.0,
                "y": 26.7,
                "w": 80.0,
                "h": 30.0,
                "role": "text",
                "font": "title_size",
                "color": "primary",
                "bold": True,
                "align": "center",
            },
            "body": {
                "x": 15.0,
                "y": 60.0,
                "w": 70.0,
                "h": 26.7,
                "role": "text",
                "font": "body_size",
                "color": "text",
                "align": "center",
            },
        },
    },
    "quote": {
        "background": "background_alt",
        "accent_bar": None,
        "slide_number": True,
        "placeholders": {
            "quote": {
                "x": 11.3,
                "y": 20.0,
                "w": 77.5,
                "h": 46.7,
                "role": "text",
                "font": "heading_size",
                "color": "heading",
                "align": "center",
            },
            "attribution": {
                "x": 11.3,
                "y": 70.0,
                "w": 77.5,
                "h": 10.0,
                "role": "text",
                "font": "body_size",
                "color": "subtitle",
                "align": "center",
            },
        },
    },
    "closing": {
        "background": "background_alt",
        "accent_bar": None,
        "slide_number": False,
        "placeholders": {
            "title": {
                "x": 11.3,
                "y": 26.7,
                "w": 77.5,
                "h": 20.0,
                "role": "text",
                "font": "title_size",
                "color": "heading",
                "bold": True,
                "align": "center",
            },
            "subtitle": {
                "x": 11.3,
                "y": 46.7,
                "w": 77.5,
                "h": 10.7,
                "role": "text",
                "font": "subtitle_size",
                "color": "subtitle",
                "align": "center",
            },
            "contact": {
                "x": 11.3,
                "y": 60.0,
                "w": 77.5,
                "h": 6.7,
                "role": "text",
                "font": "body_size",
                "color": "primary",
                "align": "center",
            },
        },
    },
}

# Layout names that correspond to old type names (backward compat)
TYPE_TO_LAYOUT = {
    "title": "title",
    "content": "content",
    "section": "section",
    "two-column": "two-column",
    "image": "image",
    "closing": "closing",
}


def _resolve_layout_positions_from_layout(
    layout: dict, slide_width: int, slide_height: int
) -> dict[str, dict[str, int]]:
    """Convert percentage-based layout placeholders to absolute EMU positions."""
    resolved: dict[str, dict[str, int]] = {}
    for name, ph in layout["placeholders"].items():
        resolved[name] = {
            "x": int(ph["x"] / 100 * slide_width),
            "y": int(ph["y"] / 100 * slide_height),
            "w": int(ph["w"] / 100 * slide_width),
            "h": int(ph["h"] / 100 * slide_height),
        }
    return resolved


def _resolve_layout_positions(
    layout_name: str, slide_width: int, slide_height: int
) -> dict[str, dict[str, int]]:
    """Convert percentage-based layout placeholders to absolute EMU positions."""
    return _resolve_layout_positions_from_layout(
        SLIDE_LAYOUTS[layout_name], slide_width, slide_height
    )


# ============================================================================
# KEYRING CREDENTIAL STORAGE
# ============================================================================


def get_credential(key: str) -> str | None:
    """Get a credential from the system keyring."""
    return keyring.get_password(SERVICE_NAME, key)


def set_credential(key: str, value: str) -> None:
    """Store a credential in the system keyring."""
    keyring.set_password(SERVICE_NAME, key, value)


def delete_credential(key: str) -> None:
    """Delete a credential from the system keyring."""
    with contextlib.suppress(keyring.errors.PasswordDeleteError):
        keyring.delete_password(SERVICE_NAME, key)


# ============================================================================
# CONFIGURATION MANAGEMENT
# ============================================================================


def load_config(service: str) -> dict[str, Any] | None:
    """Load configuration from file."""
    config_file = CONFIG_DIR / f"{service}.yaml"
    if config_file.exists():
        with open(config_file) as f:
            return yaml.safe_load(f)
    return None


def save_config(service: str, config: dict[str, Any]) -> None:
    """Save configuration to file."""
    CONFIG_DIR.mkdir(parents=True, exist_ok=True)
    config_file = CONFIG_DIR / f"{service}.yaml"
    with open(config_file, "w") as f:
        yaml.safe_dump(config, f, default_flow_style=False)


# ============================================================================
# GOOGLE AUTHENTICATION
# ============================================================================


class AuthenticationError(Exception):
    """Exception raised for authentication errors."""

    pass


def _build_oauth_config(client_id: str, client_secret: str) -> dict[str, Any]:
    """Build OAuth client configuration dict."""
    return {
        "installed": {
            "client_id": client_id,
            "client_secret": client_secret,
            "auth_uri": "https://accounts.google.com/o/oauth2/auth",
            "token_uri": "https://oauth2.googleapis.com/token",
            "redirect_uris": ["http://localhost"],
        }
    }


def get_oauth_client_config(service: str) -> dict[str, Any]:
    """Get OAuth 2.0 client configuration from config file or environment.

    Priority:
    1. Service-specific config file (~/.config/agent-skills/{service}.yaml)
    2. Service-specific environment variables
    3. Shared Google config file (~/.config/agent-skills/google.yaml)
    4. Shared environment variables (GOOGLE_CLIENT_ID, GOOGLE_CLIENT_SECRET)
    """
    config = load_config(service)
    if config and "oauth_client" in config:
        client_id = config["oauth_client"].get("client_id")
        client_secret = config["oauth_client"].get("client_secret")
        if client_id and client_secret:
            return _build_oauth_config(client_id, client_secret)

    prefix = service.upper().replace("-", "_")
    client_id = os.environ.get(f"{prefix}_CLIENT_ID")
    client_secret = os.environ.get(f"{prefix}_CLIENT_SECRET")
    if client_id and client_secret:
        return _build_oauth_config(client_id, client_secret)

    shared_config = load_config("google")
    if shared_config and "oauth_client" in shared_config:
        client_id = shared_config["oauth_client"].get("client_id")
        client_secret = shared_config["oauth_client"].get("client_secret")
        if client_id and client_secret:
            return _build_oauth_config(client_id, client_secret)

    client_id = os.environ.get("GOOGLE_CLIENT_ID")
    client_secret = os.environ.get("GOOGLE_CLIENT_SECRET")
    if client_id and client_secret:
        return _build_oauth_config(client_id, client_secret)

    raise AuthenticationError(
        f"OAuth client credentials not found for {service}. "
        f"Options:\n"
        f"  1. Run: python google-slides.py auth setup --client-id ID --client-secret SECRET\n"
        f"  2. Set env vars: GOOGLE_SLIDES_CLIENT_ID and GOOGLE_SLIDES_CLIENT_SECRET\n"
        f"  3. Create ~/.config/agent-skills/google.yaml with oauth_client credentials\n"
        f"  4. Set env vars: GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET"
    )


def _run_oauth_flow(service: str, scopes: list[str]) -> Credentials:
    """Run OAuth browser flow and store resulting token."""
    client_config = get_oauth_client_config(service)
    flow = InstalledAppFlow.from_client_config(client_config, scopes)
    creds = flow.run_local_server(port=0)
    set_credential(f"{service}-token-json", creds.to_json())
    return creds


def get_google_credentials(service: str, scopes: list[str]) -> Credentials:
    """Get Google credentials for human-in-the-loop use cases."""
    token_json = get_credential(f"{service}-token-json")
    if token_json:
        try:
            token_data = json.loads(token_json)
            creds = Credentials.from_authorized_user_info(token_data, scopes)
            if creds and creds.valid:
                granted = set(token_data.get("scopes", []))
                requested = set(scopes)
                if granted and not requested.issubset(granted):
                    merged = list(granted | requested)
                    print(
                        "Current token lacks required scopes. "
                        "Opening browser for re-authentication...",
                        file=sys.stderr,
                    )
                    delete_credential(f"{service}-token-json")
                    return _run_oauth_flow(service, merged)
                return creds
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
                set_credential(f"{service}-token-json", creds.to_json())
                return creds
        except Exception:
            pass

    try:
        return _run_oauth_flow(service, scopes)
    except Exception as e:
        raise AuthenticationError(f"OAuth flow failed: {e}") from e


def build_slides_service(scopes: list[str] | None = None):
    """Build and return Google Slides API service."""
    if scopes is None:
        scopes = SLIDES_SCOPES_DEFAULT
    creds = get_google_credentials("google-slides", scopes)
    return build("slides", "v1", credentials=creds)


def build_drive_service(scopes: list[str] | None = None):
    """Build and return Google Drive API service."""
    if scopes is None:
        scopes = DRIVE_SCOPES_READONLY
    creds = get_google_credentials("google-slides", scopes)
    return build("drive", "v3", credentials=creds)


# ============================================================================
# GOOGLE SLIDES API ERROR HANDLING
# ============================================================================


class SlidesAPIError(Exception):
    """Exception raised for Google Slides API errors."""

    def __init__(self, message: str, status_code: int | None = None, details: Any = None):
        super().__init__(message)
        self.status_code = status_code
        self.details = details


def handle_api_error(error: HttpError) -> None:
    """Convert Google API HttpError to SlidesAPIError."""
    status_code = error.resp.status
    reason = error.resp.reason
    details = None

    try:
        error_content = json.loads(error.content.decode("utf-8"))
        details = error_content.get("error", {})
        message = details.get("message", reason)
    except Exception:
        message = reason

    if status_code == 403 and "insufficient" in message.lower():
        scope_help = (
            "\n\nInsufficient OAuth scope. This operation requires additional permissions.\n"
            "To re-authenticate with the required scopes:\n\n"
            "  1. Reset token: python google-slides.py auth reset\n"
            "  2. Re-run: python google-slides.py check\n\n"
            "For setup help, see: docs/google-oauth-setup.md\n"
        )
        message = f"{message}{scope_help}"

    raise SlidesAPIError(
        f"Google Slides API error: {message} (HTTP {status_code})",
        status_code=status_code,
        details=details,
    )


# ============================================================================
# GOOGLE SLIDES API OPERATIONS
# ============================================================================


def get_presentation(service, presentation_id: str) -> dict[str, Any]:
    """Get a presentation by ID."""
    try:
        return service.presentations().get(presentationId=presentation_id).execute()
    except HttpError as e:
        handle_api_error(e)
        return {}


def read_presentation_content(service, presentation_id: str) -> str:
    """Extract text content from all slides in a presentation."""
    presentation = get_presentation(service, presentation_id)
    slides = presentation.get("slides", [])

    output_parts = []
    for idx, slide in enumerate(slides):
        slide_text = _extract_slide_text(slide)
        if slide_text.strip():
            output_parts.append(f"--- Slide {idx + 1} ---\n{slide_text}")

    return "\n\n".join(output_parts)


def _extract_slide_text(slide: dict) -> str:
    """Extract all text from a slide's page elements."""
    text_parts = []
    for element in slide.get("pageElements", []):
        if "shape" in element:
            shape = element["shape"]
            if "text" in shape:
                text = _extract_text_from_text_content(shape["text"])
                if text.strip():
                    text_parts.append(text.strip())
        if "table" in element:
            table_text = _extract_table_text(element["table"])
            if table_text.strip():
                text_parts.append(table_text)
    return "\n".join(text_parts)


def _extract_text_from_text_content(text_content: dict) -> str:
    """Extract text from a textContent structure."""
    parts = []
    for text_elem in text_content.get("textElements", []):
        if "textRun" in text_elem:
            parts.append(text_elem["textRun"].get("content", ""))
    return "".join(parts)


def _extract_table_text(table: dict) -> str:
    """Extract text from a table element as markdown."""
    rows_text = []
    for row_idx, row in enumerate(table.get("tableRows", [])):
        cell_texts = []
        for cell in row.get("tableCells", []):
            if "text" in cell:
                cell_texts.append(_extract_text_from_text_content(cell["text"]).strip())
            else:
                cell_texts.append("")
        if cell_texts:
            rows_text.append("| " + " | ".join(cell_texts) + " |")
            if row_idx == 0:
                rows_text.append("| " + " | ".join(["---"] * len(cell_texts)) + " |")
    return "\n".join(rows_text)


def format_presentation_summary(presentation: dict[str, Any]) -> str:
    """Format a presentation for display."""
    title = presentation.get("title", "(Untitled)")
    presentation_id = presentation.get("presentationId", "(Unknown)")
    slides = presentation.get("slides", [])
    return (
        f"### {title}\n"
        f"- **Presentation ID:** {presentation_id}\n"
        f"- **Slides:** {len(slides)}\n"
        f"- **URL:** https://docs.google.com/presentation/d/{presentation_id}/edit"
    )


def format_slide_info(slide: dict[str, Any], index: int) -> str:
    """Format slide information for display."""
    slide_id = slide.get("objectId", "(Unknown)")
    layout = slide.get("slideProperties", {}).get("layoutObjectId", "(Unknown)")
    elements = slide.get("pageElements", [])
    element_counts = {"shapes": 0, "images": 0, "text": 0, "other": 0}
    for element in elements:
        if "shape" in element:
            if element["shape"].get("shapeType") == "TEXT_BOX":
                element_counts["text"] += 1
            else:
                element_counts["shapes"] += 1
        elif "image" in element:
            element_counts["images"] += 1
        else:
            element_counts["other"] += 1
    total = sum(element_counts.values())
    return (
        f"### Slide {index + 1}\n"
        f"- **ID:** {slide_id}\n"
        f"- **Layout:** {layout}\n"
        f"- **Elements:** {total} ({element_counts['text']} text, "
        f"{element_counts['shapes']} shapes, {element_counts['images']} images, "
        f"{element_counts['other']} other)"
    )


def export_presentation_as_pdf(presentation_id: str) -> bytes:
    """Export presentation as PDF using Google's native export."""
    try:
        service = build_drive_service()
        return service.files().export(fileId=presentation_id, mimeType="application/pdf").execute()
    except HttpError as e:
        handle_api_error(e)
        return b""


# ============================================================================
# SLIDE LAYOUT DETECTION
# ============================================================================


def _get_text_elements(slide: dict, slide_width: float, slide_height: float) -> list[dict]:
    """Extract text elements with position, size, and font metadata from a slide."""
    elements = []
    for elem in slide.get("pageElements", []):
        if "shape" not in elem or "text" not in elem["shape"]:
            continue

        tx = elem.get("transform", {}).get("translateX", 0)
        ty = elem.get("transform", {}).get("translateY", 0)
        w = elem.get("size", {}).get("width", {}).get("magnitude", 0)

        max_font = 0
        text_parts = []
        for te in elem["shape"]["text"].get("textElements", []):
            if "textRun" in te:
                text_parts.append(te["textRun"].get("content", ""))
                fs = te["textRun"].get("style", {}).get("fontSize", {}).get("magnitude", 0)
                max_font = max(max_font, fs)

        text = "".join(text_parts).strip()
        if not text:
            continue

        elements.append(
            {
                "x_pct": tx / slide_width * 100 if slide_width else 0,
                "y_pct": ty / slide_height * 100 if slide_height else 0,
                "w_pct": w / slide_width * 100 if slide_width else 0,
                "font": max_font,
                "text": text,
            }
        )

    return elements


def _slide_has_image(slide: dict) -> bool:
    """Check if a slide contains an image element."""
    return any("image" in elem for elem in slide.get("pageElements", []))


def detect_slide_layout(slide: dict, slide_width: float, slide_height: float) -> tuple[str, float]:
    """Detect the layout of a slide by analyzing element positions and content.

    Uses heuristic rules for broad classification, then scores against
    layout templates for fine-grained matching within each category.
    Returns (layout_name, score).
    """
    text_elems = _get_text_elements(slide, slide_width, slide_height)
    has_image = _slide_has_image(slide)

    if not text_elems:
        return "content", 0.0

    max_font = max(e["font"] for e in text_elems)
    num_text = len(text_elems)

    left_texts = [e for e in text_elems if e["x_pct"] < 45]
    right_texts = [e for e in text_elems if e["x_pct"] > 45]
    has_two_cols = len(left_texts) > 1 and len(right_texts) > 1

    if has_image:
        if has_two_cols:
            return _best_scoring_layout(
                text_elems, has_image, ["two-column", "two-column-uneven", "comparison"]
            )
        return _best_scoring_layout(text_elems, has_image, ["image", "image-with-text"])

    if num_text <= 4 and max_font >= 40 and not has_two_cols:
        large_elems = [e for e in text_elems if e["font"] >= 40]
        if large_elems:
            avg_y = sum(e["y_pct"] for e in large_elems) / len(large_elems)
            if avg_y > 50:
                score = _score_layout_match(text_elems, has_image, SLIDE_LAYOUTS["closing"])
                return "closing", score
        return _best_scoring_layout(text_elems, has_image, ["title", "title-dark"])

    if num_text <= 2 and max_font >= 30:
        return _best_scoring_layout(text_elems, has_image, ["section", "quote"])

    if has_two_cols:
        return _best_scoring_layout(
            text_elems, has_image, ["two-column", "two-column-uneven", "comparison"]
        )

    return _best_scoring_layout(
        text_elems, has_image, ["content", "content-with-icon", "content-with-graphic", "data"]
    )


# Minimum score for a layout match to be considered a good fit.
# Below this, a custom layout is emitted during download.
LAYOUT_SCORE_THRESHOLD = 15.0


def _best_scoring_layout(
    text_elems: list[dict], has_image: bool, candidates: list[str]
) -> tuple[str, float]:
    """Pick the best layout from a set of candidates by scoring."""
    best_layout = candidates[0]
    best_score = -999.0

    for name in candidates:
        if name not in SLIDE_LAYOUTS:
            continue
        score = _score_layout_match(text_elems, has_image, SLIDE_LAYOUTS[name])
        if score > best_score:
            best_score = score
            best_layout = name

    return best_layout, best_score


def _score_layout_match(
    text_elems: list[dict],
    has_image: bool,
    layout_def: dict,
) -> float:
    """Score how well text elements match a layout's placeholders."""
    placeholders = layout_def["placeholders"]
    text_phs = {k: v for k, v in placeholders.items() if v["role"] != "image"}
    image_phs = {k: v for k, v in placeholders.items() if v["role"] == "image"}

    score = 0.0

    count_diff = abs(len(text_elems) - len(text_phs))
    score -= count_diff * 10

    if has_image and image_phs:
        score += 15
    elif has_image and not image_phs:
        score -= 10
    elif not has_image and image_phs:
        score -= 5

    for elem in text_elems:
        min_dist = 100.0
        best_ph = None
        for _ph_name, ph in text_phs.items():
            dx = elem["x_pct"] - ph["x"]
            dy = elem["y_pct"] - ph["y"]
            dist = (dx * dx + dy * dy) ** 0.5
            if dist < min_dist:
                min_dist = dist
                best_ph = ph

        if best_ph is not None:
            score += max(0, 20 - min_dist)

            ph_font_key = best_ph.get("font", "body_size")
            expected_font = FONT_CONFIG.get(ph_font_key, 18)
            if elem["font"] > 0:
                font_ratio = min(elem["font"], expected_font) / max(elem["font"], expected_font)
                score += font_ratio * 5

    return score


def _build_custom_layout_from_slide(
    slide: dict, slide_width: float, slide_height: float
) -> dict[str, Any]:
    """Build a custom layout definition from a slide's actual element positions.

    Used during download when no built-in layout matches well enough.
    """
    text_elems = _get_text_elements(slide, slide_width, slide_height)
    has_image = _slide_has_image(slide)

    text_elems.sort(key=lambda e: (e["y_pct"], e["x_pct"]))

    placeholders: dict[str, dict[str, Any]] = {}
    for i, elem in enumerate(text_elems):
        if i == 0 and elem["font"] >= 24:
            name = "title"
            font_key = "title_size" if elem["font"] >= 32 else "heading_size"
            ph: dict[str, Any] = {
                "x": round(elem["x_pct"], 1),
                "y": round(elem["y_pct"], 1),
                "w": round(elem["w_pct"], 1),
                "h": 12.0,
                "role": "text",
                "font": font_key,
                "color": "heading",
                "bold": True,
            }
        elif i == 1 and elem["font"] >= 16 and elem["font"] < 28:
            name = "subtitle"
            ph = {
                "x": round(elem["x_pct"], 1),
                "y": round(elem["y_pct"], 1),
                "w": round(elem["w_pct"], 1),
                "h": 10.0,
                "role": "text",
                "font": "subtitle_size",
                "color": "subtitle",
            }
        else:
            name = f"text_{i + 1}"
            ph = {
                "x": round(elem["x_pct"], 1),
                "y": round(elem["y_pct"], 1),
                "w": round(elem["w_pct"], 1),
                "h": 15.0,
                "role": "text",
                "font": "body_size",
                "color": "text",
            }
        placeholders[name] = ph

    if has_image:
        placeholders["image"] = {
            "x": 5.0,
            "y": 20.0,
            "w": 90.0,
            "h": 70.0,
            "role": "image",
        }

    return {
        "background": "background",
        "accent_bar": None,
        "slide_number": True,
        "placeholders": placeholders,
    }


def detect_slide_type(slide: dict, slide_width: float, slide_height: float) -> str:
    """Detect the type of a slide (backward-compatible wrapper).

    Returns one of the original 6 type names: title, section, content,
    two-column, image, closing.
    """
    layout, _score = detect_slide_layout(slide, slide_width, slide_height)
    _layout_to_type = {
        "title": "title",
        "title-dark": "title",
        "section": "section",
        "content": "content",
        "content-with-icon": "content",
        "content-with-graphic": "content",
        "two-column": "two-column",
        "two-column-uneven": "two-column",
        "image": "image",
        "image-with-text": "image",
        "comparison": "two-column",
        "data": "content",
        "quote": "section",
        "closing": "closing",
    }
    return _layout_to_type.get(layout, "content")


def _extract_slide_as_markdown(
    slide: dict, slide_width: float, slide_height: float
) -> tuple[str, str, float]:
    """Extract a slide's content as structured Markdown with detected layout.

    Returns:
        Tuple of (layout_name, markdown_content, match_score).
    """
    layout, score = detect_slide_layout(slide, slide_width, slide_height)
    # Use the original type-based rendering logic, mapped from layout
    layout_to_type = {
        "title": "title",
        "title-dark": "title",
        "section": "section",
        "content": "content",
        "content-with-icon": "content",
        "content-with-graphic": "content",
        "two-column": "two-column",
        "two-column-uneven": "two-column",
        "image": "image",
        "image-with-text": "image",
        "comparison": "two-column",
        "data": "content",
        "quote": "section",
        "closing": "closing",
    }
    slide_type = layout_to_type.get(layout, "content")
    text_elems = _get_text_elements(slide, slide_width, slide_height)

    if not text_elems:
        return layout, "", score

    # Sort by vertical position, then horizontal
    text_elems.sort(key=lambda e: (e["y_pct"], e["x_pct"]))

    lines: list[str] = []

    if slide_type == "two-column":
        # Separate heading (top, large font) from column content
        heading_elems = [e for e in text_elems if e["font"] >= 28 and e["y_pct"] < 25]
        body_elems = [e for e in text_elems if e not in heading_elems]

        for h in heading_elems:
            lines.append(f"# {h['text']}")

        left = [e for e in body_elems if e["x_pct"] < 45]
        right = [e for e in body_elems if e["x_pct"] >= 45]
        left.sort(key=lambda e: e["y_pct"])
        right.sort(key=lambda e: e["y_pct"])

        lines.append("")
        lines.append("<!-- left -->")
        for e in left:
            if e["font"] >= 16:
                lines.append(f"## {e['text']}")
            else:
                for bullet_line in e["text"].split("\n"):
                    bullet_line = bullet_line.strip()
                    if bullet_line:
                        lines.append(f"- {bullet_line}")

        lines.append("")
        lines.append("<!-- right -->")
        for e in right:
            if e["font"] >= 16:
                lines.append(f"## {e['text']}")
            else:
                for bullet_line in e["text"].split("\n"):
                    bullet_line = bullet_line.strip()
                    if bullet_line:
                        lines.append(f"- {bullet_line}")

    elif slide_type in ("title", "closing"):
        for e in text_elems:
            if e["font"] >= 40:
                lines.append(f"# {e['text']}")
            elif e["font"] >= 18:
                lines.append(f"## {e['text']}")
            elif e["text"]:
                lines.append(e["text"])

    else:
        # content, section, image
        for e in text_elems:
            if e["font"] >= 28:
                lines.append(f"# {e['text']}")
            elif e["font"] >= 18:
                lines.append(f"## {e['text']}")
            else:
                for bullet_line in e["text"].split("\n"):
                    bullet_line = bullet_line.strip()
                    if bullet_line:
                        lines.append(f"- {bullet_line}")

    return layout, "\n".join(lines), score


def upload_pptx_to_google(
    pptx_path: str, title: str | None = None, folder_id: str | None = None
) -> dict[str, Any]:
    """Upload a .pptx file to Google Drive, converting to native Google Slides.

    Returns:
        Dict with presentationId and URL.
    """
    from googleapiclient.http import MediaFileUpload

    service = build_drive_service(DRIVE_SCOPES_WRITE)

    file_metadata: dict[str, Any] = {
        "name": title or Path(pptx_path).stem,
        "mimeType": "application/vnd.google-apps.presentation",
    }
    if folder_id:
        file_metadata["parents"] = [folder_id]

    media = MediaFileUpload(
        pptx_path,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=True,
    )

    try:
        file = service.files().create(body=file_metadata, media_body=media, fields="id").execute()
        presentation_id = file.get("id")
        return {
            "presentationId": presentation_id,
            "url": f"https://docs.google.com/presentation/d/{presentation_id}/edit",
        }
    except HttpError as e:
        handle_api_error(e)
        return {}


def export_presentation_as_pptx(presentation_id: str, output_path: str) -> Path:
    """Download a Google Slides presentation as .pptx."""
    service = build_drive_service()
    try:
        content = (
            service.files()
            .export(
                fileId=presentation_id,
                mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            .execute()
        )
        out = Path(output_path)
        out.write_bytes(content)
        return out
    except HttpError as e:
        handle_api_error(e)
        return Path(output_path)


def extract_images_from_pptx(pptx_path: str, output_dir: str) -> dict[int, list[str]]:
    """Extract embedded images from a .pptx file.

    Returns:
        Dict mapping slide number (1-based) to list of saved image paths.
    """
    prs = Presentation(pptx_path)
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)

    images: dict[int, list[str]] = {}
    for slide_idx, slide in enumerate(prs.slides, 1):
        img_num = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                img_num += 1
                ext = shape.image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                filename = f"slide{slide_idx}_img{img_num}.{ext}"
                img_path = out / filename
                img_path.write_bytes(shape.image.blob)
                images.setdefault(slide_idx, []).append(str(img_path))

    return images


def presentation_to_markdown(
    service, presentation_id: str, image_map: dict[int, list[str]] | None = None
) -> str:
    """Convert a Google Slides presentation to Markdown with frontmatter.

    Uses layout detection to add appropriate `<!-- layout: ... -->` directives
    and structure the content (headings, bullets, columns) based on the
    spatial arrangement of text boxes.

    Args:
        service: Slides API service.
        presentation_id: The presentation ID.
        image_map: Optional dict mapping slide numbers to image file paths
            (from extract_images_from_pptx).
    """
    presentation = get_presentation(service, presentation_id)
    title = presentation.get("title", "Untitled")
    slides = presentation.get("slides", [])

    page_size = presentation.get("pageSize", {})
    slide_width = page_size.get("width", {}).get("magnitude", 9144000)
    slide_height = page_size.get("height", {}).get("magnitude", 6858000)

    # First pass: detect layouts and collect custom layout defs for low-score slides
    slide_results: list[tuple[str, str, float]] = []
    custom_layouts: dict[str, dict] = {}

    for idx, slide in enumerate(slides):
        layout, md_content, score = _extract_slide_as_markdown(slide, slide_width, slide_height)

        if score < LAYOUT_SCORE_THRESHOLD:
            custom_name = f"custom-slide-{idx + 1}"
            custom_layouts[custom_name] = _build_custom_layout_from_slide(
                slide, slide_width, slide_height
            )
            slide_results.append((custom_name, md_content, score))
        else:
            slide_results.append((layout, md_content, score))

    # Build frontmatter
    def _yaml_val(v: Any) -> str:
        if v is None:
            return "null"
        if isinstance(v, bool):
            return "true" if v else "false"
        if isinstance(v, str) and (":" in v or '"' in v or "'" in v or v.startswith("#")):
            return f'"{v}"'
        return str(v)

    fm_lines = [
        "---",
        f"title: {_yaml_val(title)}",
        f"presentation_id: {presentation_id}",
        "palette: red-hat",
        "aspect_ratio: widescreen",
    ]
    if custom_layouts:
        fm_lines.append("custom_layouts:")
        for cl_name, cl_def in custom_layouts.items():
            fm_lines.append(f"  {cl_name}:")
            fm_lines.append(f"    background: {cl_def['background']}")
            fm_lines.append(f"    accent_bar: {_yaml_val(cl_def['accent_bar'])}")
            fm_lines.append(f"    slide_number: {_yaml_val(cl_def['slide_number'])}")
            fm_lines.append("    placeholders:")
            for ph_name, ph in cl_def["placeholders"].items():
                fm_lines.append(f"      {ph_name}:")
                for k, v in ph.items():
                    fm_lines.append(f"        {k}: {_yaml_val(v)}")
    fm_lines.append("---")
    fm_lines.append("")

    lines = fm_lines

    for idx, (layout, md_content, _score) in enumerate(slide_results):
        slide_num = idx + 1
        if idx > 0:
            lines.append("---")
            lines.append("")

        if layout != "content":
            lines.append(f"<!-- layout: {layout} -->")

        if md_content.strip():
            lines.append(md_content)

        if image_map and slide_num in image_map:
            for img_path in image_map[slide_num]:
                img_name = Path(img_path).name
                lines.append("")
                lines.append(f"![]({{img_dir}}/{img_name})")

        lines.append("")

    return "\n".join(lines)


def update_presentation_replace(
    presentation_id: str,
    pptx_path: str,
) -> dict[str, Any]:
    """Replace an existing Google Slides presentation's content.

    Uploads the .pptx as a new version of the same Drive file,
    preserving the file ID, URL, sharing settings, and comments.
    """
    from googleapiclient.http import MediaFileUpload

    drive_service = build_drive_service(DRIVE_SCOPES_WRITE)

    media = MediaFileUpload(
        pptx_path,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        resumable=True,
    )

    try:
        drive_service.files().update(
            fileId=presentation_id,
            media_body=media,
        ).execute()

        return {
            "presentationId": presentation_id,
            "url": f"https://docs.google.com/presentation/d/{presentation_id}/edit",
            "mode": "replace",
        }
    except HttpError as e:
        handle_api_error(e)
        return {}


def update_presentation_append(
    presentation_id: str,
    pptx_path: str,
    position: int | None = None,
) -> dict[str, Any]:
    """Append or insert slides from a .pptx into an existing presentation.

    Uploads .pptx as a temp presentation, then uses the Slides API
    importSlides request to copy slides into the target.
    """
    temp_result = upload_pptx_to_google(pptx_path, title="_temp_update")
    temp_id = temp_result.get("presentationId")
    if not temp_id:
        raise SlidesAPIError("Failed to upload temporary presentation")

    try:
        slides_service = build_slides_service(SLIDES_SCOPES)

        temp_pres = get_presentation(slides_service, temp_id)
        temp_slides = temp_pres.get("slides", [])
        slide_ids = [s["objectId"] for s in temp_slides]

        import_request: dict[str, Any] = {
            "importSlides": {
                "presentationId": temp_id,
                "slideObjectIds": slide_ids,
            }
        }
        if position is not None:
            import_request["importSlides"]["insertionIndex"] = position

        try:
            slides_service.presentations().batchUpdate(
                presentationId=presentation_id,
                body={"requests": [import_request]},
            ).execute()
        except HttpError as e:
            handle_api_error(e)

        mode = "insert" if position is not None else "append"
        return {
            "presentationId": presentation_id,
            "url": f"https://docs.google.com/presentation/d/{presentation_id}/edit",
            "slides_added": len(slide_ids),
            "mode": mode,
        }
    finally:
        try:
            drive_service = build_drive_service(DRIVE_SCOPES_WRITE)
            drive_service.files().delete(fileId=temp_id).execute()
        except Exception:
            pass


# ============================================================================
# MARKDOWN PARSER
# ============================================================================


def parse_frontmatter(content: str) -> tuple[dict[str, Any], str]:
    """Extract YAML frontmatter from Markdown content.

    Returns:
        Tuple of (frontmatter dict, remaining content).
    """
    if not content.startswith("---"):
        return {}, content

    end_match = re.search(r"\n---\s*\n", content[3:])
    if not end_match:
        return {}, content

    fm_text = content[3 : 3 + end_match.start()]
    body = content[3 + end_match.end() :]

    try:
        frontmatter = yaml.safe_load(fm_text) or {}
    except yaml.YAMLError:
        return {}, content

    return frontmatter, body


def parse_slide_type(text: str) -> str:
    """Extract slide layout from HTML comment directive.

    Accepts both <!-- layout: name --> and <!-- type: name --> for
    backward compatibility. Returns the layout name.
    """
    match = re.search(r"<!--\s*(?:layout|type):\s*([\w-]+)\s*-->", text)
    if match:
        name = match.group(1)
        return TYPE_TO_LAYOUT.get(name, name)
    return "content"


def parse_slides(body: str) -> list[dict[str, Any]]:
    """Split Markdown body into individual slide specs."""
    raw_slides = re.split(r"\n---\s*\n", body)

    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        slides.append(_parse_single_slide(raw))

    return slides


def _parse_single_slide(raw: str) -> dict[str, Any]:
    """Parse a single slide's Markdown content into a structured dict."""
    layout = parse_slide_type(raw)
    content = re.sub(r"<!--\s*(?:layout|type):\s*[\w-]+\s*-->\s*\n?", "", raw).strip()

    slide: dict[str, Any] = {"type": layout, "layout": layout}

    if layout in ("two-column", "two-column-uneven", "comparison"):
        _parse_two_column(content, slide)
    else:
        _parse_standard_slide(content, slide)

    return slide


def _parse_standard_slide(content: str, slide: dict[str, Any]) -> None:
    """Parse a standard slide (title, section, content, image, closing, quote, data)."""
    layout = slide.get("layout", "content")
    lines = content.split("\n")
    bullets: list[str] = []
    notes: list[str] = []
    in_notes = False

    for line in lines:
        stripped = line.strip()

        if stripped.startswith("# ") and "title" not in slide:
            slide["title"] = stripped[2:].strip()
        elif stripped.startswith("## ") and "subtitle" not in slide:
            slide["subtitle"] = stripped[3:].strip()
        elif stripped.startswith("### ") and "attribution" not in slide and layout == "quote":
            slide["attribution"] = stripped[4:].strip()
        elif stripped.startswith("> "):
            if layout == "quote" and "quote" not in slide:
                slide["quote"] = stripped[2:].strip()
            else:
                notes.append(stripped[2:].strip())
                in_notes = True
        elif in_notes and stripped.startswith("> "):
            notes.append(stripped[2:].strip())
        else:
            in_notes = False
            img_match = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", stripped)
            if img_match:
                slide["image_alt"] = img_match.group(1)
                slide["image_path"] = img_match.group(2)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                bullet_text = stripped[2:].strip()
                bullet_entry = _parse_bullet_icons(bullet_text)
                bullets.append(bullet_entry)
            elif stripped and not stripped.startswith("#"):
                if layout == "closing" and "contact" not in slide:
                    slide["contact"] = stripped
                elif layout == "data" and "metric" not in slide:
                    slide["metric"] = stripped
                elif layout in ("title", "title-dark", "section"):
                    slide.setdefault("description", stripped)

    if bullets:
        slide["bullets"] = bullets
    if notes:
        slide["notes"] = "\n".join(notes)


def _parse_two_column(content: str, slide: dict[str, Any]) -> None:
    """Parse a two-column slide."""
    title_match = re.match(r"#\s+(.+)\n", content)
    if title_match:
        slide["title"] = title_match.group(1).strip()
        content = content[title_match.end() :]

    notes: list[str] = []
    clean_lines = []
    for line in content.split("\n"):
        if line.strip().startswith("> "):
            notes.append(line.strip()[2:].strip())
        else:
            clean_lines.append(line)
    if notes:
        slide["notes"] = "\n".join(notes)
    content = "\n".join(clean_lines)

    parts = re.split(r"<!--\s*(left|right)\s*-->\s*\n?", content)

    left_content = ""
    right_content = ""
    current = None
    for part in parts:
        if part.strip() == "left":
            current = "left"
        elif part.strip() == "right":
            current = "right"
        elif current == "left":
            left_content = part
        elif current == "right":
            right_content = part

    slide["left"] = _parse_column(left_content)
    slide["right"] = _parse_column(right_content)


def _parse_bullet_icons(text: str) -> dict[str, Any] | str:
    """Parse a bullet for inline icon references.

    If the bullet contains ::icon:name::, returns a dict with 'icon' and 'text'.
    Otherwise returns the plain string.
    """
    match = re.match(r"::icon:([^:]+)::\s*(.*)", text)
    if match:
        return {"icon": match.group(1).strip(), "text": match.group(2).strip()}
    return text


def _parse_column(content: str) -> dict[str, Any]:
    """Parse a column's content into heading and bullets."""
    col: dict[str, Any] = {}
    bullets: list[Any] = []
    for line in content.strip().split("\n"):
        stripped = line.strip()
        if stripped.startswith("## "):
            col["heading"] = stripped[3:].strip()
        elif stripped.startswith("- ") or stripped.startswith("* "):
            bullets.append(_parse_bullet_icons(stripped[2:].strip()))
    if bullets:
        col["bullets"] = bullets
    return col


def parse_markdown(content: str) -> dict[str, Any]:
    """Parse a full Markdown presentation file into a structured spec."""
    frontmatter, body = parse_frontmatter(content)
    slides = parse_slides(body)

    return {
        "title": frontmatter.get("title", "Untitled Presentation"),
        "palette": frontmatter.get("palette", "red-hat"),
        "aspect_ratio": frontmatter.get("aspect_ratio", "widescreen"),
        "presentation_id": frontmatter.get("presentation_id"),
        "logo_path": frontmatter.get("logo_path"),
        "custom_layouts": frontmatter.get("custom_layouts", {}),
        "slides": slides,
    }


# ============================================================================
# LIBREOFFICE INTEGRATION
# ============================================================================


def find_libreoffice() -> str | None:
    """Find LibreOffice binary on the system.

    Checks: LIBREOFFICE_PATH env var, libreoffice, soffice,
    flatpak org.libreoffice.LibreOffice.
    """
    env_path = os.environ.get("LIBREOFFICE_PATH")
    if env_path and shutil.which(env_path):
        return env_path

    for binary in ("libreoffice", "soffice"):
        path = shutil.which(binary)
        if path:
            return path

    if shutil.which("flatpak"):
        try:
            result = subprocess.run(
                ["flatpak", "info", "org.libreoffice.LibreOffice"],
                capture_output=True,
                timeout=5,
            )
            if result.returncode == 0:
                return "flatpak run org.libreoffice.LibreOffice"
        except (subprocess.TimeoutExpired, OSError):
            pass

    return None


def render_to_images(pptx_path: str, output_dir: str) -> list[str]:
    """Render .pptx slides to PNG images using LibreOffice headless."""
    lo_binary = find_libreoffice()
    if not lo_binary:
        checked = ["libreoffice", "soffice", "flatpak run org.libreoffice.LibreOffice"]
        raise RuntimeError(
            f"LibreOffice not found. Checked: {', '.join(checked)}. "
            f"Set LIBREOFFICE_PATH env var to provide the path."
        )

    Path(output_dir).mkdir(parents=True, exist_ok=True)

    cmd = lo_binary.split() + [
        "--headless",
        "--convert-to",
        "png",
        "--outdir",
        output_dir,
        pptx_path,
    ]
    subprocess.run(cmd, capture_output=True, check=True, timeout=120)

    return sorted(str(p) for p in Path(output_dir).glob("*.png"))


# ============================================================================
# ICON SYSTEM
# ============================================================================

REDHAT_ICONS_REPO = "https://github.com/RedHat-UX/red-hat-icons.git"
REDHAT_ICON_CATEGORIES = ("standard", "ui", "social", "microns")


def ensure_icon_repo(cache_dir: Path | None = None) -> Path:
    """Ensure the Red Hat icons repo is cloned locally.

    Does a shallow clone on first use; subsequent calls are instant.

    Returns:
        Path to the cloned repo's src/ directory.
    """
    cache = cache_dir or ICON_CACHE_DIR
    repo_dir = cache / "red-hat-icons"
    src_dir = repo_dir / "src"

    if src_dir.exists():
        return src_dir

    cache.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        ["git", "clone", "--depth", "1", REDHAT_ICONS_REPO, str(repo_dir)],
        capture_output=True,
        check=True,
        timeout=60,
    )
    return src_dir


def fetch_icon(name: str, cache_dir: Path | None = None) -> Path | None:
    """Find a Red Hat icon by name from the cached repo.

    Searches all categories (standard, ui, social, microns) for the icon.
    If name is a path to a local file, returns it directly.
    Clones the icon repo on first use if not already cached.

    Returns:
        Path to the SVG file, or None if not found.
    """
    icon_path = Path(name)
    if icon_path.exists():
        return icon_path

    try:
        src_dir = ensure_icon_repo(cache_dir)
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired, OSError):
        return None

    for category in REDHAT_ICON_CATEGORIES:
        svg_path = src_dir / category / f"{name}.svg"
        if svg_path.exists():
            return svg_path

    return None


def convert_svg_to_png(
    svg_path: str | Path, output_path: str | Path | None = None, scale: int = 2
) -> Path:
    """Convert an SVG file to PNG using cairosvg.

    Args:
        svg_path: Path to the SVG file.
        output_path: Output PNG path. If None, uses same name with .png extension.
        scale: Scale factor for the output (default: 2 for high-DPI).

    Returns:
        Path to the generated PNG file.
    """
    svg_path = Path(svg_path)
    output_path = svg_path.with_suffix(".png") if output_path is None else Path(output_path)

    cairosvg.svg2png(url=str(svg_path), write_to=str(output_path), scale=scale)
    return output_path


def resolve_icon(name: str, cache_dir: Path | None = None) -> Path | None:
    """Resolve an icon name to a PNG file, fetching and converting as needed.

    Handles both named Red Hat icons (e.g., "openshift") and local paths.
    SVG files are automatically converted to PNG for python-pptx compatibility.

    Returns:
        Path to a PNG file, or None if the icon couldn't be resolved.
    """
    if not CAIROSVG_AVAILABLE:
        return None

    icon_path = Path(name)

    # Local file reference
    if icon_path.exists():
        if icon_path.suffix.lower() == ".svg":
            png_path = icon_path.with_suffix(".png")
            if not png_path.exists():
                convert_svg_to_png(icon_path, png_path)
            return png_path
        return icon_path if icon_path.suffix.lower() in (".png", ".jpg", ".jpeg") else None

    # Named icon — fetch from Red Hat repo
    svg_path = fetch_icon(name, cache_dir)
    if svg_path is None:
        return None

    cache = cache_dir or ICON_CACHE_DIR
    png_path = cache / f"{name}.png"
    if not png_path.exists():
        convert_svg_to_png(svg_path, png_path)
    return png_path


# ============================================================================
# PRESENTATION BUILDER
# ============================================================================


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


TEMPLATE_DIR = Path(__file__).parent.parent / "templates"

# Placeholder styling: (idx, font_size_pt, color_palette_key, bold)
_T = FONT_CONFIG["title_size"]
_H = FONT_CONFIG["heading_size"]
_S = FONT_CONFIG["subtitle_size"]
_B = FONT_CONFIG["body_size"]
_C = FONT_CONFIG["caption_size"]

TEMPLATE_LAYOUTS: dict[str, dict[str, Any]] = {
    "title": {
        "index": 0,
        "placeholders": {
            "title": (0, _T, "heading", True),
            "subtitle": (1, _S, "subtitle", False),
        },
    },
    "title-dark": {
        "index": 1,
        "placeholders": {
            "title": (0, _T, "heading", True),
            "subtitle": (1, _S, "subtitle", False),
        },
    },
    "section": {
        "index": 2,
        "placeholders": {
            "title": (0, _T, "heading", True),
            "subtitle": (1, _S, "subtitle", False),
        },
    },
    "content": {
        "index": 3,
        "placeholders": {
            "title": (0, _H, "heading", True),
            "body": (1, _B, "text", False),
        },
    },
    "content-with-icon": {
        "index": 4,
        "placeholders": {
            "title": (0, _H, "heading", True),
            "body": (1, _B, "text", False),
            "icon": (10, _B, "text", False),
        },
    },
    "content-with-graphic": {
        "index": 5,
        "placeholders": {
            "title": (0, _H, "heading", True),
            "body": (1, _B, "text", False),
            "graphic": (10, _B, "text", False),
        },
    },
    "two-column": {
        "index": 6,
        "placeholders": {
            "title": (0, _H, "heading", True),
            "left_heading": (10, _S, "primary", True),
            "left_body": (11, _B, "text", False),
            "right_heading": (12, _S, "primary", True),
            "right_body": (13, _B, "text", False),
        },
    },
    "two-column-uneven": {
        "index": 7,
        "placeholders": {
            "title": (0, _H, "heading", True),
            "left_heading": (10, _S, "primary", True),
            "left_body": (11, _B, "text", False),
            "right_heading": (12, _S, "primary", True),
            "right_body": (13, _B, "text", False),
        },
    },
    "image": {
        "index": 8,
        "placeholders": {
            "title": (0, _H, "heading", True),
            "image": (10, _B, "text", False),
            "caption": (11, _C, "subtitle", False),
        },
    },
    "image-with-text": {
        "index": 9,
        "placeholders": {
            "title": (0, _H, "heading", True),
            "image": (10, _B, "text", False),
            "body": (1, _B, "text", False),
        },
    },
    "comparison": {
        "index": 10,
        "placeholders": {
            "title": (0, _H, "heading", True),
            "left_heading": (10, _S, "accent", True),
            "left_body": (11, _B, "text", False),
            "right_heading": (12, _S, "primary", True),
            "right_body": (13, _B, "text", False),
        },
    },
    "data": {
        "index": 11,
        "placeholders": {
            "title": (0, _H, "heading", True),
            "metric": (10, _T, "primary", True),
            "body": (1, _B, "text", False),
        },
    },
    "quote": {
        "index": 12,
        "placeholders": {
            "quote": (0, _H, "heading", False),
            "attribution": (1, _B, "subtitle", False),
        },
    },
    "closing": {
        "index": 13,
        "placeholders": {
            "title": (0, _T, "heading", True),
            "subtitle": (1, _S, "subtitle", False),
            "contact": (10, _B, "primary", False),
        },
    },
}


class PresentationBuilder:
    """Builds professional .pptx presentations from a template with slide layouts."""

    def __init__(self, palette_name: str = "red-hat", aspect_ratio: str = "widescreen"):
        if palette_name not in COLOR_PALETTES:
            raise ValueError(
                f"Unknown palette '{palette_name}'. Available: {', '.join(sorted(COLOR_PALETTES))}"
            )

        self.palette = COLOR_PALETTES[palette_name]
        self.palette_name = palette_name
        self.aspect_ratio = aspect_ratio

        template_path = TEMPLATE_DIR / "default.pptx"
        self.prs = Presentation(str(template_path))

        if aspect_ratio == "widescreen":
            self.prs.slide_width = SLIDE_WIDTH_WIDESCREEN
            self.prs.slide_height = SLIDE_HEIGHT_WIDESCREEN
        else:
            self.prs.slide_width = SLIDE_WIDTH_STANDARD
            self.prs.slide_height = SLIDE_HEIGHT_STANDARD

        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

    def build_from_spec(self, spec: dict[str, Any], output_path: str | None = None) -> Path:
        """Build a presentation from a parsed spec dict."""
        self._custom_layouts = spec.get("custom_layouts", {})

        for slide_spec in spec.get("slides", []):
            layout = slide_spec.get("layout", slide_spec.get("type", "content"))
            if layout in TYPE_TO_LAYOUT:
                layout = TYPE_TO_LAYOUT[layout]
            if layout not in TEMPLATE_LAYOUTS and layout not in self._custom_layouts:
                layout = "content"
            self._add_layout_slide(layout, slide_spec)

        if output_path:
            out = Path(output_path)
        else:
            out = Path(f"{spec.get('title', 'presentation').replace(' ', '_')}.pptx")

        self.prs.save(str(out))
        return out

    def _add_speaker_notes(self, slide, notes: str) -> None:
        """Add speaker notes to a slide."""
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    @staticmethod
    def _suppress_bullets(p) -> None:
        """Explicitly disable bullets on a paragraph."""
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        pPr = p._p.get_or_add_pPr()
        for old in pPr.findall(f"{{{ns}}}buChar") + pPr.findall(f"{{{ns}}}buAutoNum"):
            pPr.remove(old)
        if pPr.find(f"{{{ns}}}buNone") is None:
            pPr.append(pPr.makeelement(f"{{{ns}}}buNone", {}))

    @staticmethod
    def _set_bullet_format(p) -> None:
        """Explicitly enable bullet character with indent on a paragraph."""
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(Inches(0.35)))
        pPr.set("indent", str(-Inches(0.25)))
        for old in pPr.findall(f"{{{ns}}}buNone"):
            pPr.remove(old)
        if pPr.find(f"{{{ns}}}buChar") is None:
            pPr.append(pPr.makeelement(f"{{{ns}}}buFont", {"typeface": "Arial", "charset": "0"}))
            pPr.append(pPr.makeelement(f"{{{ns}}}buChar", {"char": "•"}))

    def _fill_text_placeholder(
        self,
        slide,
        ph_idx: int,
        text: str,
        font_size: int = 20,
        color_key: str = "text",
        bold: bool = False,
    ) -> None:
        """Fill a text placeholder with explicit font styling and no bullets."""
        try:
            placeholder = slide.placeholders[ph_idx]
        except KeyError:
            return

        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        self._suppress_bullets(p)
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = hex_to_rgb(self.palette.get(color_key, self.palette["text"]))
            run.font.name = FONT_CONFIG["font_family"]

    def _fill_bullet_placeholder(
        self,
        slide,
        ph_idx: int,
        bullets: list[Any],
        font_size: int = 20,
        color_key: str = "text",
    ) -> None:
        """Fill a bullet placeholder with explicit bullet and font styling."""
        try:
            placeholder = slide.placeholders[ph_idx]
        except KeyError:
            return

        tf = placeholder.text_frame
        tf.clear()

        for i, bullet in enumerate(bullets):
            if isinstance(bullet, dict):
                text = bullet.get("text", "")
                icon_name = bullet.get("icon", "")
            else:
                text = str(bullet)
                icon_name = ""

            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.level = 0
            self._set_bullet_format(p)
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.color.rgb = hex_to_rgb(self.palette.get(color_key, self.palette["text"]))
                run.font.name = FONT_CONFIG["font_family"]

            if icon_name:
                icon_path = resolve_icon(icon_name)
                if icon_path and icon_path.exists():
                    slide.shapes.add_picture(
                        str(icon_path),
                        placeholder.left,
                        placeholder.top + i * (Pt(FONT_CONFIG["body_size"]) + Pt(8)),
                        Inches(0.35),
                        Inches(0.35),
                    )

    def _fill_image_placeholder(self, slide, ph_idx: int, image_path: str) -> None:
        """Fill a picture placeholder with an image."""
        if not image_path or not Path(image_path).exists():
            return
        try:
            placeholder = slide.placeholders[ph_idx]
            placeholder.insert_picture(image_path)
        except (KeyError, AttributeError):
            pass

    def _add_layout_slide(self, layout_name: str, spec: dict[str, Any]) -> None:
        """Add a slide using a template layout and fill its placeholders."""
        if layout_name not in TEMPLATE_LAYOUTS:
            layout_name = "content"

        tmpl = TEMPLATE_LAYOUTS[layout_name]
        slide_layout = self.prs.slide_layouts[tmpl["index"]]
        slide = self.prs.slides.add_slide(slide_layout)

        content = self._map_spec_to_placeholders(layout_name, spec)
        ph_map = tmpl["placeholders"]

        for content_name, ph_info in ph_map.items():
            ph_idx, font_size, color_key, bold = ph_info
            data = content.get(content_name)
            if not data:
                continue

            if isinstance(data, list):
                self._fill_bullet_placeholder(
                    slide, ph_idx, data, font_size=font_size, color_key=color_key
                )
            elif content_name in ("image", "icon", "graphic"):
                self._fill_image_placeholder(slide, ph_idx, data)
            else:
                self._fill_text_placeholder(
                    slide,
                    ph_idx,
                    str(data),
                    font_size=font_size,
                    color_key=color_key,
                    bold=bold,
                )

        if spec.get("notes"):
            self._add_speaker_notes(slide, spec["notes"])

    def _map_spec_to_placeholders(self, _layout_name: str, spec: dict[str, Any]) -> dict[str, Any]:
        """Map slide spec data to layout placeholder names."""
        content: dict[str, Any] = {}

        content["title"] = spec.get("title", "")
        subtitle = spec.get("subtitle", "")
        description = spec.get("description", "")
        if description and subtitle:
            subtitle = f"{subtitle}\n{description}"
        elif description:
            subtitle = description
        content["subtitle"] = subtitle
        content["description"] = description
        content["body"] = spec.get("bullets", []) or []
        content["contact"] = spec.get("contact", "")
        content["caption"] = spec.get("image_alt") or spec.get("caption", "")

        # Quote layout
        content["quote"] = spec.get("quote", "")
        content["attribution"] = spec.get("attribution", "")

        # Data layout
        content["metric"] = spec.get("metric", "")

        # Image placeholders
        image_path = spec.get("image_path", "")
        content["image"] = image_path
        content["icon"] = spec.get("icon_path", "")
        content["graphic"] = spec.get("graphic_path") or image_path

        # Two-column data
        left = spec.get("left", {})
        right = spec.get("right", {})
        content["left_heading"] = left.get("heading", "")
        content["left_body"] = left.get("bullets", []) or []
        content["right_heading"] = right.get("heading", "")
        content["right_body"] = right.get("bullets", []) or []

        return content


# ============================================================================
# HEALTH CHECK
# ============================================================================


def check_slides_connectivity() -> dict[str, Any]:
    """Check Google Slides API connectivity and authentication."""
    result: dict[str, Any] = {
        "authenticated": False,
        "scopes": None,
        "error": None,
    }

    try:
        creds = get_google_credentials("google-slides", SLIDES_SCOPES_DEFAULT)
        available_scopes: list[str] = []
        if hasattr(creds, "scopes"):
            available_scopes = creds.scopes
        elif hasattr(creds, "_scopes"):
            available_scopes = creds._scopes

        service = build("slides", "v1", credentials=creds)
        test_pres = service.presentations().create(body={"title": "_test_connectivity"}).execute()
        test_pres_id = test_pres.get("presentationId")

        result["authenticated"] = True
        result["test_presentation_id"] = test_pres_id
        result["scopes"] = {
            "readonly": any("presentations.readonly" in s for s in available_scopes),
            "write": any("presentations" in s and "readonly" not in s for s in available_scopes),
            "all_scopes": available_scopes,
        }
    except Exception as e:
        result["error"] = str(e)

    return result


# ============================================================================
# CLI COMMAND HANDLERS
# ============================================================================


def cmd_check(_args):
    """Handle 'check' command."""
    print("Checking Google Slides connectivity...")
    result = check_slides_connectivity()

    if result["authenticated"]:
        print("✓ Successfully authenticated to Google Slides")
        scopes = result.get("scopes", {})
        if scopes:
            print("\nGranted OAuth Scopes:")
            print(f"  Read-only (presentations.readonly): {'✓' if scopes.get('readonly') else '✗'}")
            print(f"  Write (presentations):               {'✓' if scopes.get('write') else '✗'}")
            if not scopes.get("write"):
                print("\n⚠️  Write scope not granted. Some operations will fail.")
                print("   To grant full access, reset and re-authenticate:")
                print()
                print("   1. Reset token: python google-slides.py auth reset")
                print("   2. Re-run: python google-slides.py check")
                print()
                print("   See: docs/google-oauth-setup.md")
        print(f"\nTest presentation created: {result.get('test_presentation_id')}")
        print("(You can delete this test presentation from Google Drive)")
        return 0
    else:
        print(f"✗ Authentication failed: {result['error']}")
        print()
        print("Setup instructions:")
        print()
        print("  1. Set up a GCP project with OAuth credentials:")
        print("     See: docs/gcp-project-setup.md")
        print()
        print("  2. Configure your credentials:")
        print("     Create ~/.config/agent-skills/google.yaml:")
        print()
        print("     oauth_client:")
        print("       client_id: YOUR_CLIENT_ID.apps.googleusercontent.com")
        print("       client_secret: YOUR_CLIENT_SECRET")
        print()
        print("  3. Run check again to trigger OAuth flow:")
        print("     python google-slides.py check")
        print()
        print("For detailed setup instructions, see: docs/google-oauth-setup.md")
        return 1


def cmd_auth_setup(args):
    """Handle 'auth setup' command."""
    if not args.client_id or not args.client_secret:
        print("Error: Both --client-id and --client-secret are required", file=sys.stderr)
        return 1

    config = load_config("google-slides") or {}
    config["oauth_client"] = {
        "client_id": args.client_id,
        "client_secret": args.client_secret,
    }
    save_config("google-slides", config)
    print("✓ OAuth client credentials saved to config file")
    print(f"  Config location: {CONFIG_DIR / 'google-slides.yaml'}")
    print("\nNext step: Run any Google Slides command to initiate OAuth flow")
    return 0


def cmd_auth_reset(_args):
    """Handle 'auth reset' command."""
    delete_credential("google-slides-token-json")
    print("OAuth token cleared. Next command will trigger re-authentication.")
    return 0


def cmd_auth_status(_args):
    """Handle 'auth status' command."""
    token_json = get_credential("google-slides-token-json")
    if not token_json:
        print("No OAuth token stored.")
        return 1

    try:
        token_data = json.loads(token_json)
    except json.JSONDecodeError:
        print("Stored token is corrupted.")
        return 1

    print("OAuth token is stored.")
    scopes = token_data.get("scopes", [])
    if scopes:
        print("\nGranted scopes:")
        for scope in scopes:
            print(f"  - {scope}")
    else:
        print("\nGranted scopes: (unknown - legacy token)")

    has_refresh = bool(token_data.get("refresh_token"))
    print(f"\nRefresh token: {'present' if has_refresh else 'missing'}")

    expiry = token_data.get("expiry")
    if expiry:
        print(f"Token expiry: {expiry}")

    client_id = token_data.get("client_id", "")
    if client_id:
        truncated = client_id[:16] + "..." if len(client_id) > 16 else client_id
        print(f"Client ID: {truncated}")

    return 0


def cmd_create(args):
    """Handle 'create' command — build .pptx from Markdown, optionally upload."""
    if not args.file:
        print("Error: --file is required for create", file=sys.stderr)
        return 1

    if not PPTX_AVAILABLE:
        print(
            "Error: python-pptx not found. Install with: pip install --user python-pptx",
            file=sys.stderr,
        )
        return 1

    file_path = Path(args.file)
    if not file_path.exists():
        print(f"Error: File not found: {args.file}", file=sys.stderr)
        return 1

    markdown_content = file_path.read_text()
    spec = parse_markdown(markdown_content)

    palette = args.palette or spec.get("palette", "red-hat")

    try:
        builder = PresentationBuilder(
            palette_name=palette, aspect_ratio=spec.get("aspect_ratio", "widescreen")
        )
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    output_path = args.output or str(file_path.with_suffix(".pptx"))

    # Resolve image paths relative to the Markdown file
    md_dir = file_path.parent
    for slide_spec in spec.get("slides", []):
        if "image_path" in slide_spec:
            img = Path(slide_spec["image_path"])
            if not img.is_absolute():
                slide_spec["image_path"] = str(md_dir / img)

    result_path = builder.build_from_spec(spec, output_path=output_path)

    print(f"✓ Built: {result_path}")
    print(f"  Slides: {len(spec.get('slides', []))}")
    print(f"  Palette: {palette}")

    if args.title:
        try:
            upload_result = upload_pptx_to_google(str(result_path), title=args.title)
            print("✓ Uploaded to Google Slides")
            print(f"  Presentation ID: {upload_result['presentationId']}")
            print(f"  URL: {upload_result['url']}")
        except (SlidesAPIError, AuthenticationError) as e:
            print(f"Upload failed: {e}", file=sys.stderr)
            return 1

    return 0


def cmd_get(args):
    """Handle 'get' command — get presentation content or download as .pptx."""
    service = build_slides_service(SLIDES_SCOPES_READONLY)
    presentation = get_presentation(service, args.presentation_id)

    if args.output:
        pptx_path = Path(args.output)
        if pptx_path.suffix != ".pptx":
            pptx_path = pptx_path.with_suffix(".pptx")
        export_presentation_as_pptx(args.presentation_id, str(pptx_path))
        print(f"✓ Saved PPTX: {pptx_path}")
        return 0

    if args.json:
        print(json.dumps(presentation, indent=2))
    else:
        print(format_presentation_summary(presentation))
        print()
        content = read_presentation_content(service, args.presentation_id)
        if content:
            print(content)

    return 0


def cmd_update(args):
    """Handle 'update' command — update existing presentation from local .pptx."""
    if not Path(args.file).exists():
        print(f"Error: File not found: {args.file}", file=sys.stderr)
        return 1

    try:
        if args.mode == "replace":
            result = update_presentation_replace(args.presentation_id, args.file)
        else:
            position = args.position if args.mode == "insert" else None
            result = update_presentation_append(args.presentation_id, args.file, position)

        print(f"✓ Updated presentation ({result['mode']} mode)")
        print(f"  URL: {result['url']}")
        return 0
    except (SlidesAPIError, AuthenticationError) as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1


def cmd_preview(args):
    """Handle 'preview' command — render .pptx to images via LibreOffice."""
    if not Path(args.file).exists():
        print(f"Error: File not found: {args.file}", file=sys.stderr)
        return 1

    fmt = args.format or "images"
    if fmt == "summary":
        if not PPTX_AVAILABLE:
            print("Error: python-pptx not found.", file=sys.stderr)
            return 1
        prs = Presentation(args.file)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            print(f"--- Slide {i + 1} ---")
            for text in texts:
                print(f"  {text}")
            print()
        return 0
    else:
        try:
            output_dir = str(Path(args.file).with_suffix("")) + "_preview"
            images = render_to_images(args.file, output_dir)
            print(f"✓ Rendered {len(images)} slide(s) to: {output_dir}")
            for img in images:
                print(f"  {img}")
            return 0
        except RuntimeError as e:
            print(f"Error: {e}", file=sys.stderr)
            return 1


def _image_dimensions_from_blob(blob: bytes) -> tuple[int | None, int | None]:
    """Extract pixel dimensions from PNG or JPEG image blob without PIL."""
    if blob[:8] == b"\x89PNG\r\n\x1a\n" and len(blob) >= 24:
        w = int.from_bytes(blob[16:20], "big")
        h = int.from_bytes(blob[20:24], "big")
        return w, h

    if blob[:2] == b"\xff\xd8":
        i = 2
        while i < len(blob) - 9:
            if blob[i] != 0xFF:
                break
            marker = blob[i + 1]
            if marker in (0xC0, 0xC1, 0xC2):
                h = int.from_bytes(blob[i + 5 : i + 7], "big")
                w = int.from_bytes(blob[i + 7 : i + 9], "big")
                return w, h
            seg_len = int.from_bytes(blob[i + 2 : i + 4], "big")
            i += 2 + seg_len

    return None, None


def verify_presentation(pptx_path: str) -> list[dict[str, Any]]:
    """Run quality and accessibility checks on a local .pptx file.

    Returns:
        List of check results, each with keys: check, severity, message.
    """
    prs = Presentation(pptx_path)
    results: list[dict[str, Any]] = []

    # Slide count
    slide_count = len(prs.slides)
    results.append({"check": "slide_count", "severity": "info", "message": f"{slide_count} slides"})

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        has_title = False
        fonts_used: set[str] = set()

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
                    for run in para.runs:
                        if run.font.name:
                            fonts_used.add(run.font.name)
                        if run.font.size and run.font.size >= Pt(24):
                            has_title = True

        # Empty slide check
        if not slide_texts:
            results.append(
                {
                    "check": "empty_slide",
                    "severity": "warning",
                    "message": f"Slide {i}: no text content",
                }
            )

        # Title presence check
        if not has_title and i > 1:
            results.append(
                {
                    "check": "title_presence",
                    "severity": "warning",
                    "message": f"Slide {i}: no title-sized text (≥24pt) found",
                }
            )

        # Font consistency check
        expected_font = FONT_CONFIG["font_family"]
        unexpected = fonts_used - {expected_font}
        if unexpected:
            results.append(
                {
                    "check": "font_consistency",
                    "severity": "warning",
                    "message": f"Slide {i}: unexpected fonts: {', '.join(sorted(unexpected))}",
                }
            )

    # Color contrast check (WCAG AA)
    for i, slide in enumerate(prs.slides, 1):
        bg_color = None
        try:
            fill = slide.background.fill
            if fill.type is not None and fill.fore_color and fill.fore_color.rgb:
                bg_color = f"#{fill.fore_color.rgb}"
        except (AttributeError, TypeError):
            pass

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_bg = None
            try:
                sfill = shape.fill
                if sfill.type is not None and sfill.fore_color and sfill.fore_color.rgb:
                    shape_bg = f"#{sfill.fore_color.rgb}"
            except (AttributeError, TypeError):
                pass
            effective_bg = shape_bg or bg_color or "#FFFFFF"

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    fg_color = None
                    try:
                        if run.font.color.rgb:
                            fg_color = f"#{run.font.color.rgb}"
                    except (AttributeError, TypeError):
                        pass
                    if not fg_color:
                        continue

                    ratio = contrast_ratio(fg_color, effective_bg)
                    font_size = run.font.size
                    is_bold = run.font.bold
                    is_large = (font_size and font_size >= Pt(18)) or (
                        is_bold and font_size and font_size >= Pt(14)
                    )
                    threshold = 3.0 if is_large else 4.5

                    if ratio < threshold:
                        results.append(
                            {
                                "check": "color_contrast",
                                "severity": "warning",
                                "message": (
                                    f"Slide {i}: low contrast {ratio:.1f}:1 "
                                    f"({fg_color} on {effective_bg}), "
                                    f"WCAG AA requires {threshold:.1f}:1"
                                ),
                            }
                        )

    # Text overflow detection
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.height or not shape.width:
                continue

            total_height_emu = 0
            for para in shape.text_frame.paragraphs:
                text = para.text
                if not text:
                    total_height_emu += Pt(8)
                    continue

                font_size = Pt(FONT_CONFIG["body_size"])
                for run in para.runs:
                    if run.font.size:
                        font_size = run.font.size
                        break

                avg_char_width = int(font_size * 0.6)
                chars_per_line = max(1, shape.width // avg_char_width) if avg_char_width else 1
                line_count = max(1, -(-len(text) // chars_per_line))
                total_height_emu += line_count * int(font_size * 1.3) + Pt(8)

            if total_height_emu > shape.height:
                results.append(
                    {
                        "check": "text_overflow",
                        "severity": "info",
                        "message": (
                            f"Slide {i}: text may overflow its container "
                            f"(estimated {total_height_emu / EMU_PER_PT:.0f}pt "
                            f"in {shape.height / EMU_PER_PT:.0f}pt box)"
                        ),
                    }
                )

    # Image resolution check (150 DPI minimum)
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not hasattr(shape, "image"):
                continue
            try:
                blob = shape.image.blob
                pixel_w, pixel_h = _image_dimensions_from_blob(blob)
                if pixel_w and pixel_h and shape.width and shape.height:
                    dpi_x = pixel_w / (shape.width / EMU_PER_INCH)
                    dpi_y = pixel_h / (shape.height / EMU_PER_INCH)
                    min_dpi = min(dpi_x, dpi_y)
                    if min_dpi < 150:
                        results.append(
                            {
                                "check": "image_dpi",
                                "severity": "warning",
                                "message": (
                                    f"Slide {i}: image resolution {min_dpi:.0f} DPI "
                                    f"(minimum 150 DPI recommended)"
                                ),
                            }
                        )
            except (AttributeError, ValueError):
                pass

    # Speaker notes check
    slides_with_notes = []
    slides_without_notes = []
    for i, slide in enumerate(prs.slides, 1):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            slides_with_notes.append(i)
        else:
            slides_without_notes.append(i)

    if slides_without_notes:
        results.append(
            {
                "check": "speaker_notes",
                "severity": "info",
                "message": f"Slides without notes: {', '.join(map(str, slides_without_notes))}",
            }
        )

    return results


def relative_luminance(hex_color: str) -> float:
    """Calculate relative luminance of a color per WCAG 2.0."""
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16) / 255
    g = int(hex_color[2:4], 16) / 255
    b = int(hex_color[4:6], 16) / 255

    def linearize(c: float) -> float:
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    return 0.2126 * linearize(r) + 0.7152 * linearize(g) + 0.0722 * linearize(b)


def contrast_ratio(color1: str, color2: str) -> float:
    """Calculate WCAG contrast ratio between two hex colors."""
    l1 = relative_luminance(color1)
    l2 = relative_luminance(color2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def cmd_verify(args):
    """Handle 'verify' command — quality checks on .pptx."""
    if not Path(args.file).exists():
        print(f"Error: File not found: {args.file}", file=sys.stderr)
        return 1

    if not PPTX_AVAILABLE:
        print("Error: python-pptx not found.", file=sys.stderr)
        return 1

    if args.presentation_id:
        # Compare local vs cloud
        service = build_slides_service(SLIDES_SCOPES_READONLY)
        cloud_pres = get_presentation(service, args.presentation_id)
        cloud_slides = cloud_pres.get("slides", [])

        prs = Presentation(args.file)
        local_slide_count = len(prs.slides)
        cloud_slide_count = len(cloud_slides)

        print(f"\nCloud comparison: {args.presentation_id}")
        if local_slide_count == cloud_slide_count:
            print(f"  ✓ Slide count matches: {local_slide_count}")
        else:
            print(f"  ⚠ Slide count mismatch: local={local_slide_count}, cloud={cloud_slide_count}")

        # Compare text content
        cloud_content = read_presentation_content(service, args.presentation_id)
        local_texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            local_texts.append(text)
        local_content = "\n".join(local_texts)

        if local_content.strip() and cloud_content.strip():
            local_words = set(local_content.lower().split())
            cloud_words = set(cloud_content.lower().split())
            overlap = len(local_words & cloud_words)
            total = len(local_words | cloud_words)
            similarity = overlap / total if total > 0 else 0
            print(f"  Text similarity: {similarity:.0%}")

        return 0

    results = verify_presentation(args.file)

    warnings = [r for r in results if r["severity"] == "warning"]
    infos = [r for r in results if r["severity"] == "info"]

    print(f"Verification: {args.file}\n")

    if infos:
        for r in infos:
            print(f"  ℹ {r['message']}")

    if warnings:
        print()
        for r in warnings:
            print(f"  ⚠ {r['message']}")

    print(f"\n{len(warnings)} warning(s), {len(infos)} info(s)")
    return 0


def cmd_palettes(_args):
    """Handle 'palettes' command — list available color palettes."""
    print("Available color palettes:\n")
    for name, colors in sorted(COLOR_PALETTES.items()):
        marker = " (default)" if name == "red-hat" else ""
        print(f"### {name}{marker}")
        print(f"  Primary:    {colors['primary']}")
        print(f"  Background: {colors['background']}")
        print(f"  Text:       {colors['text']}")
        print(f"  Accent:     {colors['accent']}")
        print(f"  Divider:    {colors['divider']}")
        print()
    return 0


# ============================================================================
# CLI ARGUMENT PARSER
# ============================================================================


def build_parser() -> argparse.ArgumentParser:
    """Build the argument parser."""
    parser = argparse.ArgumentParser(
        description="Google Slides — Markdown-driven presentation builder",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )

    subparsers = parser.add_subparsers(dest="command", help="Command to execute")

    # check command
    subparsers.add_parser("check", help="Check Google Slides connectivity and authentication")

    # auth commands
    auth_parser = subparsers.add_parser("auth", help="Authentication management")
    auth_subparsers = auth_parser.add_subparsers(dest="auth_command")
    setup_parser = auth_subparsers.add_parser("setup", help="Setup OAuth client credentials")
    setup_parser.add_argument("--client-id", required=True, help="OAuth client ID")
    setup_parser.add_argument("--client-secret", required=True, help="OAuth client secret")
    auth_subparsers.add_parser("reset", help="Clear stored OAuth token")
    auth_subparsers.add_parser("status", help="Show current token info")

    # create command
    create_parser = subparsers.add_parser(
        "create", help="Build .pptx from Markdown, optionally upload to Google Slides"
    )
    create_parser.add_argument("--file", required=True, help="Markdown file to build from")
    create_parser.add_argument("--output", "-o", help="Output .pptx path (default: <file>.pptx)")
    create_parser.add_argument("--palette", help="Color palette (default: red-hat)")
    create_parser.add_argument(
        "--title", help="If provided, also upload to Google Slides with this title"
    )

    # get command
    get_parser = subparsers.add_parser("get", help="Get presentation content or download as .pptx")
    get_parser.add_argument("presentation_id", help="Presentation ID")
    get_parser.add_argument("--output", "-o", help="Download as .pptx to this path")
    get_parser.add_argument("--json", action="store_true", help="Output as JSON")

    # update command
    update_parser = subparsers.add_parser(
        "update", help="Update existing Google Slides presentation from local .pptx"
    )
    update_parser.add_argument("presentation_id", help="Presentation ID")
    update_parser.add_argument("--file", required=True, help="Local .pptx file")
    update_parser.add_argument(
        "--mode",
        choices=["replace", "append", "insert"],
        default="replace",
        help="Update mode (default: replace)",
    )
    update_parser.add_argument("--position", type=int, help="Insert position (for insert mode)")

    # preview command
    preview_parser = subparsers.add_parser(
        "preview", help="Render .pptx slides to images via LibreOffice"
    )
    preview_parser.add_argument("--file", required=True, help="Local .pptx file")
    preview_parser.add_argument(
        "--format",
        choices=["images", "summary"],
        default="images",
        help="Preview format (default: images)",
    )

    # verify command
    verify_parser = subparsers.add_parser(
        "verify", help="Run quality/accessibility checks on .pptx"
    )
    verify_parser.add_argument("--file", required=True, help="Local .pptx file")
    verify_parser.add_argument(
        "--presentation-id", help="Compare with uploaded Google Slides version"
    )

    # palettes command
    subparsers.add_parser("palettes", help="List available color palettes")

    return parser


# ============================================================================
# MAIN
# ============================================================================


def main():
    """Main entry point."""
    parser = build_parser()
    args = parser.parse_args()

    if not YAML_AVAILABLE:
        print(
            "Error: 'pyyaml' library not found. Install with: pip install --user pyyaml",
            file=sys.stderr,
        )
        return 1

    if not args.command:
        parser.print_help()
        return 1

    # Commands that don't require Google API dependencies
    local_commands = {"create", "preview", "verify", "palettes"}

    if args.command not in local_commands:
        if not GOOGLE_AUTH_AVAILABLE:
            print(
                "Error: Google auth libraries not found. Install with: "
                "pip install --user google-auth google-auth-oauthlib",
                file=sys.stderr,
            )
            return 1

        if not GOOGLE_API_CLIENT_AVAILABLE:
            print(
                "Error: 'google-api-python-client' not found. Install with: "
                "pip install --user google-api-python-client",
                file=sys.stderr,
            )
            return 1

        if not KEYRING_AVAILABLE:
            print(
                "Error: 'keyring' library not found. Install with: pip install --user keyring",
                file=sys.stderr,
            )
            return 1

    try:
        if args.command == "check":
            return cmd_check(args)
        elif args.command == "auth":
            if args.auth_command == "setup":
                return cmd_auth_setup(args)
            elif args.auth_command == "reset":
                return cmd_auth_reset(args)
            elif args.auth_command == "status":
                return cmd_auth_status(args)
        elif args.command == "create":
            return cmd_create(args)
        elif args.command == "get":
            return cmd_get(args)
        elif args.command == "update":
            return cmd_update(args)
        elif args.command == "preview":
            return cmd_preview(args)
        elif args.command == "verify":
            return cmd_verify(args)
        elif args.command == "palettes":
            return cmd_palettes(args)

        parser.print_help()
        return 1

    except (SlidesAPIError, AuthenticationError) as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    except KeyboardInterrupt:
        print("\nInterrupted", file=sys.stderr)
        return 130
    except Exception as e:
        print(f"Unexpected error: {e}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
